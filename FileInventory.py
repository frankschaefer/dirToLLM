import os
import pathlib
import requests
import pdfplumber
import docx
from pptx import Presentation
from openpyxl import load_workbook
import time
import json
from datetime import datetime
import select
import sys
import warnings
import argparse

# Unterdrücke openpyxl Warnungen für nicht unterstützte Excel-Features
warnings.filterwarnings('ignore', category=UserWarning, module='openpyxl')

# Version und Metadaten
VERSION = "1.18.0"
VERSION_DATE = "2025-12-30"
SCRIPT_NAME = "FileInventory - OneDrive Dokumenten-Zusammenfassung (macOS)"

# Fehlerbehandlungsmodus: None = fragen, "skip" = weiter ohne Fragen, "ask" = weiter mit Fragen
ERROR_HANDLING_MODE = None

# macOS Pfade - expandiere ~ zum Home-Verzeichnis
SRC_ROOT = os.path.expanduser("~/OneDrive - CompanyName")
DST_ROOT = os.path.expanduser("~/LLM")

LMSTUDIO_API_URL = "http://localhost:1234/v1/chat/completions"
MODEL_NAME = "local-model"  # in LM Studio unter Model-Name des laufenden Servers schauen

# Alternativ: Falls LM Studio auf einem anderen Port läuft:
# LMSTUDIO_API_URL = "http://localhost:8080/v1/chat/completions"
# oder prüfen Sie in LM Studio unter "Local Server" welcher Port verwendet wird

# Minimale Dateigröße für Bilddateien (in Bytes) - ignoriere kleine Icons
MIN_IMAGE_SIZE = 10 * 1024  # 10 KB

# Modell Context-Länge (maximale Anzahl Tokens)
# Passen Sie dies an Ihr Modell an:
# - Kleinere Modelle (z.B. Llama 3 8B): 8192
# - Größere Modelle (z.B. Qwen 2.5 14B): 32768
# - Reasoning-Modelle (z.B. mistralai/ministral-3-14b-reasoning): 262144
MAX_CONTEXT_TOKENS = 262144

# Maximale Länge der Zusammenfassung in Zeichen
# Wenn der Originaltext kürzer ist, wird er direkt kopiert
SUMMARY_MAX_CHARS = 1500

# Welche Dateitypen sollen verarbeitet werden?
EXTENSIONS = {
    ".pdf",                                    # PDF-Dokumente
    ".docx", ".doc",                          # Word-Dokumente (neu und alt)
    ".pptx", ".ppt",                          # PowerPoint-Präsentationen (neu und alt)
    ".xlsx", ".xls", ".xlsm", ".xltx",       # Excel-Dateien (neu, alt, Makro, Vorlagen)
    ".txt", ".md",                            # Textdateien
    ".png", ".jpg", ".jpeg"                   # Bilddateien
}

# Verzeichnismuster die übersprungen werden sollen (Glob-Patterns)
EXCLUDE_PATTERNS = [
    "**/Vorlagen/**",
    "**/Templates/**",
    "**/Musterdateien/**",
    "**/1) Musterdateien und Vorlagen/**",
    "**/_archive/**",
    "**/.old/**",
    "**/backup/**",
    "**/Backup/**",
]

# Duplikat-Erkennung: Cache für Dateigrößen und Hashes
_SIZE_HASH_CACHE = {}  # {size: {hash: path}}

# ============================================================================
# DSGVO / BDSG - Klassifizierung besonders schutzbedürftiger Daten
# ============================================================================
# Gemäß Art. 9 DSGVO (besondere Kategorien personenbezogener Daten)
# und § 26 BDSG (Beschäftigtendaten)

SENSITIVE_DATA_KEYWORDS = {
    # § 26 BDSG - Beschäftigtendaten
    "GEHALTSABRECHNUNG": {
        "keywords": ["lohnabrechnung", "gehaltsabrechnung", "entgeltabrechnung",
                    "bruttolohn", "nettolohn", "bruttogehalt", "nettogehalt",
                    "lohnsteuer", "sozialversicherung", "entgelt",
                    "verdienst", "lohnzettel", "gehaltsmitteilung"],
        "dsgvo_kategorie": "Art. 9 Abs. 2 lit. b DSGVO i.V.m. § 26 BDSG - Beschäftigtendaten",
        "schutzklasse": "hoch"
    },
    "LEBENSLAUF": {
        "keywords": ["lebenslauf", "curriculum vitae", "cv", "bewerbung", "werdegang",
                    "beruflicher werdegang", "vita", "bewerbungsunterlagen", "qualifikation"],
        "dsgvo_kategorie": "§ 26 BDSG - Beschäftigtendaten (Bewerber)",
        "schutzklasse": "hoch"
    },
    "ARBEITSVERTRAG": {
        "keywords": ["arbeitsvertrag", "anstellungsvertrag", "dienstvertrag", "arbeitgeber",
                    "arbeitnehmer", "arbeitsverh\u00e4ltnis", "vertragspartei", "probezeit",
                    "k\u00fcndigung", "befristet", "unbefristet"],
        "dsgvo_kategorie": "§ 26 BDSG - Beschäftigtendaten",
        "schutzklasse": "hoch"
    },
    "ZEUGNIS": {
        "keywords": ["arbeitszeugnis", "zwischenzeugnis", "zeugnis", "beurteilung",
                    "leistungsbeurteilung", "qualifiziertes zeugnis", "einfaches zeugnis"],
        "dsgvo_kategorie": "§ 26 BDSG - Beschäftigtendaten",
        "schutzklasse": "hoch"
    },
    "PERSONALAKTE": {
        "keywords": ["personalakte", "personaldaten", "mitarbeiterdaten", "personalstammdaten",
                    "personalnummer", "mitarbeiter", "besch\u00e4ftigte"],
        "dsgvo_kategorie": "§ 26 BDSG - Beschäftigtendaten",
        "schutzklasse": "sehr hoch"
    },

    # Art. 9 DSGVO - Gesundheitsdaten
    "GESUNDHEITSDATEN": {
        "keywords": ["attest", "arbeitsunf\u00e4higkeit", "krankheit", "arzt", "gesundheit",
                    "schwerbehinderung", "au-bescheinigung", "krankmeldung", "medizinisch",
                    "diagnose", "therapie", "reha", "betriebsarzt"],
        "dsgvo_kategorie": "Art. 9 Abs. 1 DSGVO - Gesundheitsdaten",
        "schutzklasse": "sehr hoch"
    },

    # Sozialversicherung und Steuern (§ 26 BDSG)
    "SOZIALVERSICHERUNG": {
        "keywords": ["sozialversicherungsnummer", "rentenversicherung", "krankenversicherung",
                    "sv-nummer", "versicherungsnummer", "krankenkasse", "rentenversicherungsnummer"],
        "dsgvo_kategorie": "§ 26 BDSG - Beschäftigtendaten (Sozialversicherung)",
        "schutzklasse": "sehr hoch"
    },
    "STEUER": {
        "keywords": ["lohnsteuerbescheinigung", "steuernummer", "finanzamt", "steuerklasse",
                    "steuer-id", "identifikationsnummer", "elstam", "lohnsteuer"],
        "dsgvo_kategorie": "§ 26 BDSG - Beschäftigtendaten (Steuerdaten)",
        "schutzklasse": "sehr hoch"
    },

    # Ausweisdokumente
    "AUSWEIS": {
        "keywords": ["personalausweis", "reisepass", "ausweisnummer", "pass-nummer",
                    "ausweiskopie", "identit\u00e4tsnachweis", "ausweisdokument"],
        "dsgvo_kategorie": "Art. 6 Abs. 1 DSGVO - Identifikationsdaten",
        "schutzklasse": "sehr hoch"
    },


    # BANKDATEN wurde ENTFERNT - Begründung:
    # - Firmen-IBANs sind NICHT personenbezogen (Art. 4 Nr. 1 DSGVO)
    # - Nur Private Bankverbindungen natürlicher Personen sind schützenswert
    # - Wird nun via LLM-Kontext-Check geprüft (siehe classify_sensitive_data)
}

# Prüfe OCR-Verfügbarkeit global (einmalig beim Start)
OCR_AVAILABLE = False
pytesseract = None
PIL_Image = None
try:
    import pytesseract
    from PIL import Image as PIL_Image
    OCR_AVAILABLE = True
except ImportError:
    pass  # OCR nicht verfügbar

def extract_text_pdf(path):
    """
    Extrahiert Text aus PDF-Dateien.
    Verwendet OCR (Tesseract) für gescannte PDFs ohne Text.

    Returns:
        tuple: (text, ocr_info) wobei ocr_info ein dict ist mit:
            - 'used_ocr': Boolean, ob OCR verwendet wurde
            - 'ocr_pages': Anzahl der Seiten mit OCR
            - 'total_pages': Gesamtzahl der Seiten
            - 'ocr_chars': Anzahl der via OCR extrahierten Zeichen
    """
    texts = []
    ocr_pages = 0
    total_ocr_chars = 0
    ocr_info = {
        'used_ocr': False,
        'ocr_pages': 0,
        'total_pages': 0,
        'ocr_chars': 0
    }

    try:
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)
            ocr_info['total_pages'] = total_pages

            # Verarbeite jede Seite
            for page_num, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text() or ""

                # Wenn keine oder sehr wenig Text gefunden wurde, könnte es ein Scan sein
                if len(page_text.strip()) < 10:
                    # Versuche OCR mit pytesseract (falls verfügbar)
                    if OCR_AVAILABLE:
                        try:
                            # Konvertiere PDF-Seite zu Bild
                            if hasattr(page, 'to_image'):
                                pil_image = page.to_image(resolution=300).original

                                # OCR mit Tesseract (Deutsch)
                                ocr_text = pytesseract.image_to_string(pil_image, lang='deu')

                                if len(ocr_text.strip()) > len(page_text.strip()):
                                    page_text = ocr_text
                                    ocr_pages += 1
                                    total_ocr_chars += len(ocr_text)
                                    ocr_info['used_ocr'] = True

                                    if page_num == 1:
                                        print(f"  → OCR verwendet für Seite {page_num}/{total_pages}")

                        except Exception as e:
                            # OCR fehlgeschlagen, verwende ursprünglichen Text
                            if page_num == 1:
                                print(f"  → OCR-Fehler auf Seite {page_num}: {str(e)[:50]}")
                    else:
                        # pytesseract nicht installiert - nur einmal warnen
                        if page_num == 1:
                            print(f"  → Warnung: OCR nicht verfügbar (pytesseract nicht installiert)")

                texts.append(page_text)

                # Zeige Fortschritt bei vielen Seiten
                if total_pages > 10 and page_num % 10 == 0:
                    print(f"  → PDF-Verarbeitung: {page_num}/{total_pages} Seiten")

    except Exception as e:
        print(f"  → Fehler beim PDF-Öffnen: {e}")
        return "", ocr_info

    result = "\n\n".join(texts)

    # Update OCR Info
    ocr_info['ocr_pages'] = ocr_pages
    ocr_info['ocr_chars'] = total_ocr_chars

    if ocr_info['used_ocr'] and len(result.strip()) > 100:
        print(f"  → OCR Ergebnis: {ocr_pages}/{total_pages} Seiten mit OCR verarbeitet, {total_ocr_chars:,} Zeichen extrahiert")

    return result, ocr_info

def extract_text_docx(path):
    """Extrahiert Text aus Word-Dokumenten (.docx)."""
    doc = docx.Document(path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_doc(path):
    """
    Extrahiert Text aus alten Word-Dokumenten (.doc).
    Hinweis: .doc-Format wird nicht nativ unterstützt.
    Als Workaround wird versucht, mit python-docx zu öffnen (funktioniert manchmal).
    """
    try:
        # python-docx kann manchmal auch .doc öffnen (wenn es eigentlich .docx ist)
        doc = docx.Document(path)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        # Fallback: Rückgabe mit Hinweis
        return f"[.doc-Datei - Textextraktion nicht vollständig möglich. Benötigt LibreOffice/antiword für vollständige Konvertierung]"

def extract_text_pptx(path):
    """Extrahiert Text aus PowerPoint-Dateien (.pptx)."""
    texts = []
    try:
        prs = Presentation(path)
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if slide_texts:
                texts.append(f"Folie {slide_num}:\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)
    except Exception as e:
        print(f"Warnung bei PPTX-Extraktion: {e}")
        return ""

def extract_text_ppt(path):
    """
    Extrahiert Text aus alten PowerPoint-Dateien (.ppt).
    Hinweis: .ppt-Format wird nicht nativ unterstützt.
    Als Workaround wird versucht, mit python-pptx zu öffnen (funktioniert manchmal).
    """
    try:
        # python-pptx kann manchmal auch .ppt öffnen (wenn es eigentlich .pptx ist)
        prs = Presentation(path)
        texts = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text)
            if slide_texts:
                texts.append(f"Folie {slide_num}:\n" + "\n".join(slide_texts))
        return "\n\n".join(texts)
    except Exception as e:
        # Fallback: Rückgabe mit Hinweis
        return f"[.ppt-Datei - Textextraktion nicht vollständig möglich. Benötigt LibreOffice für vollständige Konvertierung]"

def extract_text_xlsx(path):
    """Extrahiert Text (keine Formeln) aus Excel-Dateien (.xlsx, .xlsm, .xltx)."""
    texts = []
    try:
        wb = load_workbook(path, data_only=True)  # data_only=True gibt Werte statt Formeln
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts = []

            for row in sheet.iter_rows(values_only=True):
                # Filtere None-Werte und konvertiere zu String
                row_texts = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if row_texts:
                    sheet_texts.append(" | ".join(row_texts))

            if sheet_texts:
                texts.append(f"Arbeitsblatt '{sheet_name}':\n" + "\n".join(sheet_texts))

        wb.close()
        return "\n\n".join(texts)
    except Exception as e:
        print(f"Warnung bei Excel-Extraktion: {e}")
        return ""

def extract_text_xls(path):
    """
    Extrahiert Text aus alten Excel-Dateien (.xls).
    Hinweis: .xls-Format wird von openpyxl nicht unterstützt.
    Benötigt xlrd-Bibliothek für vollständige Unterstützung.
    """
    try:
        # Versuche mit openpyxl (funktioniert nur wenn Datei fälschlicherweise .xls heißt)
        wb = load_workbook(path, data_only=True)
        texts = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheet_texts = []
            for row in sheet.iter_rows(values_only=True):
                row_texts = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
                if row_texts:
                    sheet_texts.append(" | ".join(row_texts))
            if sheet_texts:
                texts.append(f"Arbeitsblatt '{sheet_name}':\n" + "\n".join(sheet_texts))
        wb.close()
        return "\n\n".join(texts)
    except Exception as e:
        # Fallback: Rückgabe mit Hinweis
        return f"[.xls-Datei - Textextraktion nicht möglich. Benötigt xlrd-Bibliothek oder LibreOffice für Konvertierung]"

def extract_text_txt(path):
    """Extrahiert Text aus TXT-Dateien."""
    try:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        # Fallback für andere Encodings
        try:
            with open(path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            print(f"Warnung bei TXT-Extraktion: {e}")
            return ""
    except Exception as e:
        print(f"Warnung bei TXT-Extraktion: {e}")
        return ""

def extract_text_md(path):
    """Extrahiert Text aus Markdown-Dateien."""
    try:
        with open(path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        try:
            with open(path, 'r', encoding='latin-1') as f:
                return f.read()
        except Exception as e:
            print(f"Warnung bei MD-Extraktion: {e}")
            return ""
    except Exception as e:
        print(f"Warnung bei MD-Extraktion: {e}")
        return ""

def extract_text_image(path):
    """Für Bilddateien wird ein Platzhalter zurückgegeben - das Bild wird per Vision API analysiert."""
    # Der eigentliche Text wird später vom LLM extrahiert, das die Bilddatei direkt analysiert
    return f"[IMAGE_FILE:{path}]"

def extract_text(path):
    """
    Extrahiert Text aus einer Datei.

    Returns:
        tuple: (text, ocr_info) wobei ocr_info None ist für nicht-PDF Dateien
    """
    ext = path.suffix.lower()
    if ext == ".pdf":
        return extract_text_pdf(path)  # Gibt (text, ocr_info) zurück
    elif ext == ".docx":
        return extract_text_docx(path), None
    elif ext == ".doc":
        return extract_text_doc(path), None
    elif ext == ".pptx":
        return extract_text_pptx(path), None
    elif ext == ".ppt":
        return extract_text_ppt(path), None
    elif ext in {".xlsx", ".xlsm", ".xltx"}:
        return extract_text_xlsx(path), None
    elif ext == ".xls":
        return extract_text_xls(path), None
    elif ext == ".txt":
        return extract_text_txt(path), None
    elif ext == ".md":
        return extract_text_md(path), None
    elif ext in {".png", ".jpg", ".jpeg"}:
        return extract_text_image(path), None
    else:
        return "", None

def is_file_accessible(file_path):
    """
    Prüft, ob eine Datei zugänglich ist.
    Auf macOS sind OneDrive-Dateien normalerweise direkt verfügbar.
    """
    try:
        return os.path.exists(file_path) and os.access(file_path, os.R_OK)
    except Exception as e:
        print(f"Warnung: Konnte Dateizugriff nicht prüfen für {file_path}: {e}")
        return False

def should_exclude_path(path):
    """
    Prüft ob ein Pfad basierend auf EXCLUDE_PATTERNS übersprungen werden soll.

    Args:
        path: Pfad als String oder pathlib.Path

    Returns:
        True wenn Pfad ausgeschlossen werden soll, False sonst
    """
    import fnmatch

    path_str = str(path)
    rel_path = os.path.relpath(path_str, SRC_ROOT)

    for pattern in EXCLUDE_PATTERNS:
        # Nutze fnmatch für Glob-Pattern-Matching
        if fnmatch.fnmatch(rel_path, pattern.lstrip('**/')):
            return True
        # Prüfe auch absoluten Pfad
        if fnmatch.fnmatch(path_str, pattern):
            return True

    return False

def calculate_content_hash(file_path):
    """
    Berechnet SHA-256 Hash des Dateiinhalts für Duplikat-Erkennung.

    Args:
        file_path: Pfad zur Datei

    Returns:
        SHA-256 Hash als Hex-String
    """
    import hashlib

    hasher = hashlib.sha256()
    try:
        with open(file_path, 'rb') as f:
            # Lese in Chunks für große Dateien
            for chunk in iter(lambda: f.read(8192), b''):
                hasher.update(chunk)
        return hasher.hexdigest()
    except Exception as e:
        print(f"  → Warnung: Konnte Hash nicht berechnen für {file_path}: {e}")
        return None

def is_duplicate_file(file_path, file_size):
    """
    Prüft ob eine Datei ein Duplikat einer bereits verarbeiteten Datei ist.
    Verwendet Größe + Content-Hash für effiziente Duplikat-Erkennung.

    Args:
        file_path: Pfad zur zu prüfenden Datei
        file_size: Größe der Datei in Bytes

    Returns:
        (is_duplicate, original_path) - Tuple mit Boolean und Pfad zum Original (oder None)
    """
    global _SIZE_HASH_CACHE

    # Schritt 1: Prüfe ob Dateigröße bereits bekannt
    if file_size not in _SIZE_HASH_CACHE:
        # Neue Größe - kann kein Duplikat sein
        _SIZE_HASH_CACHE[file_size] = {}
        file_hash = calculate_content_hash(file_path)
        if file_hash:
            _SIZE_HASH_CACHE[file_size][file_hash] = file_path
        return False, None

    # Schritt 2: Größe existiert - berechne Hash und prüfe
    file_hash = calculate_content_hash(file_path)
    if not file_hash:
        # Hash-Berechnung fehlgeschlagen - behandle nicht als Duplikat
        return False, None

    # Schritt 3: Prüfe ob Hash bereits existiert
    if file_hash in _SIZE_HASH_CACHE[file_size]:
        original_path = _SIZE_HASH_CACHE[file_size][file_hash]
        return True, original_path

    # Schritt 4: Neuer Hash für diese Größe - speichere
    _SIZE_HASH_CACHE[file_size][file_hash] = file_path
    return False, None

def extract_contact_info_from_text(text):
    """
    Extrahiert URLs, E-Mail-Adressen und Telefonnummern aus Text mittels Regex.
    Sehr schnell (keine LLM-Aufrufe), ideal für inkrementelle Updates.

    Args:
        text: Der zu durchsuchende Text

    Returns:
        dict: {'urls': [...], 'emails': [...], 'phone_numbers': [...]}
    """
    import re

    contact_info = {
        'urls': [],
        'emails': [],
        'phone_numbers': []
    }

    if not text:
        return contact_info

    # URL-Extraktion
    # Findet http(s)://, www., und gängige Domains
    # WICHTIG: Schließt Satzzeichen am Ende aus und stoppt bei @-Zeichen (E-Mail-Grenze)
    url_pattern = r'(?:https?://|www\.)(?:[A-Za-z0-9\-._~:/?#\[\]!$&\'()*+,;=%]+[A-Za-z0-9\-_~/?#\[\]$&*+=%]|[A-Za-z0-9\-._~:/?#\[\]!$&\'()*+,;=%])'
    raw_urls = re.findall(url_pattern, text, re.IGNORECASE)

    # Bereinige URLs: Entferne trailing Satzzeichen und @ (stoppt bei E-Mail)
    cleaned_urls = []
    for url in raw_urls:
        # Stoppe bei @ (trennt URL von E-Mail)
        if '@' in url:
            url = url.split('@')[0]
        # Entferne trailing Satzzeichen: ), ., ,, ;, :, !, und -Buchstaben (vor E-Mail local part)
        url = re.sub(r'[).,;:!]+$', '', url)
        # Entferne "-Wort" Pattern am Ende (z.B. "-Hallo" vor E-Mail)
        url = re.sub(r'-[A-Za-z]+$', '', url)
        # Nur URLs mit mindestens einem . im Domain-Teil
        if '.' in url and len(url) > 5:
            cleaned_urls.append(url)

    contact_info['urls'] = list(set(cleaned_urls))  # Duplikate entfernen

    # E-Mail-Extraktion
    # WICHTIG: Striktes Pattern - nur alphanumerische Zeichen, ., _, %, +, - im local part
    # Kein www. oder andere URL-Prefixe vor dem @
    email_pattern = r'\b[A-Za-z0-9][A-Za-z0-9._%+-]*@[A-Za-z0-9][A-Za-z0-9.-]*\.[A-Za-z]{2,}\b'
    raw_emails = re.findall(email_pattern, text)

    # Bereinige E-Mails: Entferne ungültige Prefixe
    cleaned_emails = []
    for email in raw_emails:
        # Prüfe ob URL-Muster im local part (www., http)
        local_part = email.split('@')[0]
        if 'www.' in local_part.lower() or 'http' in local_part.lower():
            # Extrahiere nur den Teil nach dem letzten '-' oder Leerzeichen
            # z.B. "BOOKPLAYGmbH-www.book-play.de-Hallo@book-play.de" -> "Hallo@book-play.de"
            parts = re.split(r'[-\s]', local_part)
            if parts:
                local_part = parts[-1]
                email = f"{local_part}@{email.split('@')[1]}"

        # Validierung: local part sollte nicht zu lang sein (max 64 Zeichen)
        if len(local_part) <= 64 and len(email) < 254:
            cleaned_emails.append(email)

    contact_info['emails'] = list(set(cleaned_emails))

    # Telefonnummer-Extraktion (verschiedene Formate)
    # Deutsche Formate: +49, 0049, (0), mit/ohne Leerzeichen, Bindestriche, Klammern
    # WICHTIG: Strenge Pattern um False Positives zu vermeiden (z.B. Projektnummern)
    phone_patterns = [
        r'\+49[\s\-]?\(?\d{2,4}\)?[\s\-]?\d{3,10}',  # +49 30 12345678 oder +49(30)12345678
        r'\+49[\s\-]?\d{2,4}[\s\-/]\d{6,10}',        # +49 30/12345678
        r'0049[\s\-]?\d{2,4}[\s\-]?\d{6,10}',        # 0049 30 12345678
        r'\(0\d{2,4}\)[\s\-]?\d{6,10}',              # (030) 12345678 (Mindest 6 Ziffern nach Vorwahl!)
        r'\b0\d{2,4}[\s\-/]\d{6,10}\b',              # 030/12345678 (Mindest 6 Ziffern!)
        r'\b0\d{9,11}\b',                             # 03012345678 (ohne Separator, mind. 10 Ziffern)
    ]

    phone_numbers = []
    for pattern in phone_patterns:
        matches = re.findall(pattern, text)
        phone_numbers.extend(matches)

    # Bereinige und dedupliziere Telefonnummern
    cleaned_phones = []
    for phone in phone_numbers:
        # Normalisiere: Entferne Leerzeichen für Vergleich
        normalized = re.sub(r'[\s\-/()]', '', phone)
        digits_only = re.sub(r'\D', '', normalized)

        # Strikte Validierung:
        # - Mindestens 8 Ziffern (echte Telefonnummern)
        # - Nicht nur 4-stellige Jahreszahlen (z.B. "2024")
        # - Nicht kurze Nummern wie "091-2024" (nur 7 Ziffern ohne führende 0)
        if len(digits_only) >= 8:
            # Prüfe ob es wie eine echte Telefonnummer aussieht
            # Beginnt mit 0, +49, oder 0049?
            if digits_only.startswith('0') or digits_only.startswith('49'):
                # Behalte Originalformat für Lesbarkeit
                if phone not in cleaned_phones:
                    cleaned_phones.append(phone)

    contact_info['phone_numbers'] = cleaned_phones

    return contact_info

def extract_entities_from_path(file_path):
    """
    Extrahiert potenzielle Firmen-/Projektnamen aus dem Dateipfad.

    Analysiert die Verzeichnisstruktur und identifiziert wahrscheinliche:
    - Firmennamen (z.B. "Siemens AG", "BMW Group")
    - Projektnamen (z.B. "Projekt_Digitalisierung_2024")

    Args:
        file_path: Pfad zur Datei

    Returns:
        dict: {'companies': [...], 'projects': [...]}
    """
    import re

    entities = {
        'companies': [],
        'projects': []
    }

    # Extrahiere relative Pfad-Komponenten
    rel_path = os.path.relpath(file_path, SRC_ROOT)
    path_parts = pathlib.Path(rel_path).parts

    # Typische Firmennamen-Patterns
    company_indicators = [
        r'\b\w+\s+(AG|GmbH|SE|KG|OHG|mbH|Inc\.|Corp\.|Ltd\.|Group)\b',  # Rechtsformen
        r'\b(Firma|Company|Corporation)\s+\w+\b',
    ]

    # Typische Projekt-Patterns
    project_indicators = [
        r'\bProjekt[e]?\b',
        r'\bProject[s]?\b',
        r'\b\d{4}[-_]\d{2}\b',  # z.B. "2024-01" oder "2024_Q1"
    ]

    for part in path_parts:
        # Ignoriere bekannte System-Verzeichnisse
        if part in ['Vorlagen', 'Templates', 'Musterdateien', 'backup', 'Backup', '_archive']:
            continue

        # Prüfe auf Firmennamen-Pattern
        for pattern in company_indicators:
            matches = re.findall(pattern, part, re.IGNORECASE)
            for match in matches:
                # Extrahiere vollständigen Namen (nicht nur die Rechtsform)
                # Erweitere um umgebende Wörter
                full_match = re.search(r'\b[\w\s]+' + re.escape(match) + r'\b', part, re.IGNORECASE)
                if full_match:
                    company_name = full_match.group(0).strip()
                    if company_name not in entities['companies'] and len(company_name) > 2:
                        entities['companies'].append(company_name)

        # Prüfe auf Projekt-Pattern
        for pattern in project_indicators:
            if re.search(pattern, part, re.IGNORECASE):
                # Bereinige Unterstriche und Bindestriche für bessere Lesbarkeit
                project_name = part.replace('_', ' ').replace('-', ' ')
                if project_name not in entities['projects'] and len(project_name) > 3:
                    entities['projects'].append(project_name)
                break

    return entities

def check_bankdata_context_with_llm(text):
    """
    Prüft via LLM, ob Bankdaten (IBAN/Kontonummern) im Kontext natürlicher oder juristischer Personen stehen.

    Rechtlicher Hintergrund:
    - Firmen-IBANs (GmbH, AG, etc.) sind NICHT personenbezogen (Art. 4 Nr. 1 DSGVO)
    - Private IBANs natürlicher Personen (Einzelunternehmer, Freiberufler) sind personenbezogen

    Args:
        text: Der zu analysierende Text (sollte bereits IBAN/Kontonummer enthalten)

    Returns:
        dict: {
            'contains_private_bankdata': bool,  # True wenn natürliche Person
            'confidence': str,  # 'hoch', 'mittel', 'niedrig'
            'context': str  # Kurze Erklärung
        }
    """
    import re

    # Prüfe ob überhaupt Bankdaten enthalten sind
    bankdata_pattern = r'\b(iban|kontonummer|bankverbindung|bic)\b'
    if not re.search(bankdata_pattern, text.lower()):
        return {
            'contains_private_bankdata': False,
            'confidence': 'hoch',
            'context': 'Keine Bankdaten gefunden'
        }

    # Begrenze Text auf relevante Länge
    max_chars = 3000
    if len(text) > max_chars:
        # Nehme ersten Teil (wo meist Absender/Kontext steht)
        analysis_text = text[:max_chars]
    else:
        analysis_text = text

    prompt = """Analysiere den folgenden Text auf Bankverbindungen (IBAN, Kontonummern).

AUFGABE:
Bestimme, ob die Bankdaten zu NATÜRLICHEN PERSONEN oder JURISTISCHEN PERSONEN gehören.

RECHTLICHER HINTERGRUND:
- Firmen-IBANs (GmbH, AG, UG, KG, OHG, etc.) sind NICHT personenbezogen
- Private IBANs von Einzelpersonen, Einzelunternehmer, Freiberufler SIND personenbezogen
- Geschäftsbriefe mit Firmen-IBAN sind normal und nicht besonders schützenswert

INDIKATOREN FÜR FIRMEN (→ NICHT schützenswert):
- Rechtsformen: GmbH, AG, UG, KG, OHG, e.V., Stiftung
- Kontext: Geschäftsbrief, Angebot, Rechnung an Firma
- Briefkopf mit Firmennamen und Registernummer
- "Geschäftsführer:", "Vorstand:", "Handelsregister:"

INDIKATOREN FÜR NATÜRLICHE PERSONEN (→ schützenswert):
- Einzelperson ohne Rechtsform
- Freiberufler, Selbständige (ohne GmbH/UG)
- Private Bankverbindung im Arbeitsvertrag
- Gehaltsabrechnung, Lohnzettel
- Bewerbungsunterlagen

AUSGABEFORMAT (exakt so):
TYP: [FIRMA / NATÜRLICHE_PERSON / UNKLAR]
KONFIDENZ: [HOCH / MITTEL / NIEDRIG]
KONTEXT: [Ein Satz Begründung]

Beispiele:
- Geschäftsbrief mit "Müller GmbH, IBAN DE..." → TYP: FIRMA
- Gehaltsabrechnung mit IBAN → TYP: NATÜRLICHE_PERSON
- Rechnung von "Max Mustermann Steuerberater, IBAN..." → TYP: NATÜRLICHE_PERSON
"""

    try:
        payload = {
            "model": MODEL_NAME,
            "messages": [
                {
                    "role": "system",
                    "content": "Du bist ein DSGVO-Klassifizierungs-Experte. Analysiere sachlich und präzise."
                },
                {
                    "role": "user",
                    "content": f"{prompt}\n\nTEXT:\n{analysis_text}"
                }
            ],
            "temperature": 0.1,
            "max_tokens": 200
        }

        response = requests.post(
            LMSTUDIO_API_URL,
            headers={"Content-Type": "application/json"},
            json=payload,
            timeout=30
        )

        if response.status_code == 200:
            result_text = response.json()['choices'][0]['message']['content'].strip()

            # Parse Antwort
            is_private = False
            confidence = 'niedrig'
            context = result_text

            if 'TYP:' in result_text:
                if 'NATÜRLICHE_PERSON' in result_text or 'NATÜRLICHE PERSON' in result_text:
                    is_private = True
                elif 'FIRMA' in result_text:
                    is_private = False

            if 'KONFIDENZ:' in result_text:
                if 'HOCH' in result_text:
                    confidence = 'hoch'
                elif 'MITTEL' in result_text:
                    confidence = 'mittel'

            # Extrahiere Kontext-Zeile
            for line in result_text.split('\n'):
                if line.startswith('KONTEXT:'):
                    context = line.replace('KONTEXT:', '').strip()
                    break

            return {
                'contains_private_bankdata': is_private,
                'confidence': confidence,
                'context': context
            }
        else:
            # Fallback bei LLM-Fehler: Im Zweifel als schützenswert markieren
            return {
                'contains_private_bankdata': True,
                'confidence': 'niedrig',
                'context': 'LLM-Analyse fehlgeschlagen - im Zweifel als schützenswert markiert'
            }

    except Exception as e:
        # Fallback: Im Zweifel als schützenswert markieren
        return {
            'contains_private_bankdata': True,
            'confidence': 'niedrig',
            'context': f'Fehler bei LLM-Analyse: {str(e)[:100]}'
        }

def classify_sensitive_data(text, file_path=None):
    """
    Klassifiziert Dokumente hinsichtlich besonders schutzbedürftiger personenbezogener Daten
    gemäß Art. 9 DSGVO und § 26 BDSG.

    Analysiert Dateiname, Pfad und Textinhalt auf Schlüsselbegriffe für:
    - Beschäftigtendaten (§ 26 BDSG)
    - Gesundheitsdaten (Art. 9 DSGVO)
    - Weitere sensible personenbezogene Daten

    Args:
        text: Der zu analysierende Dokumententext
        file_path: Optional - Pfad zur Datei (für Dateiname-Analyse)

    Returns:
        dict: {
            'contains_sensitive_data': bool,
            'data_categories': [str],  # Liste erkannter Kategorien
            'dsgvo_classification': [str],  # DSGVO-Artikel
            'protection_level': str,  # 'hoch' oder 'sehr hoch'
            'matched_keywords': {kategorie: [keywords]}  # Gefundene Keywords pro Kategorie
        }
    """
    import re

    result = {
        'contains_sensitive_data': False,
        'data_categories': [],
        'dsgvo_classification': [],
        'protection_level': None,
        'matched_keywords': {}
    }

    # Kombiniere Text und Dateiname für Analyse
    search_text = text.lower()
    if file_path:
        filename = os.path.basename(file_path).lower()
        search_text = filename + " " + search_text

    highest_protection = None
    protection_levels = {'hoch': 1, 'sehr hoch': 2}

    # Prüfe jede Kategorie
    for category_name, category_data in SENSITIVE_DATA_KEYWORDS.items():
        matched_keywords = []

        # Prüfe Keywords in dieser Kategorie
        for keyword in category_data['keywords']:
            # Verwende Word-Boundary für präzise Treffer
            pattern = r'\b' + re.escape(keyword) + r'\b'
            if re.search(pattern, search_text, re.IGNORECASE):
                matched_keywords.append(keyword)

        # Falls mindestens 1 Keyword gefunden wurde
        if matched_keywords:
            result['contains_sensitive_data'] = True
            result['data_categories'].append(category_name)
            result['matched_keywords'][category_name] = matched_keywords

            # Füge DSGVO-Kategorie hinzu (ohne Duplikate)
            dsgvo_cat = category_data['dsgvo_kategorie']
            if dsgvo_cat not in result['dsgvo_classification']:
                result['dsgvo_classification'].append(dsgvo_cat)

            # Aktualisiere höchste Schutzklasse
            current_level = category_data['schutzklasse']
            if highest_protection is None or protection_levels.get(current_level, 0) > protection_levels.get(highest_protection, 0):
                highest_protection = current_level

    result['protection_level'] = highest_protection

    # ZUSÄTZLICHE PRÜFUNG: Bankdaten mit LLM-Kontext-Analyse
    # Prüfe ob Dokument IBAN/Kontonummer enthält
    bankdata_pattern = r'\b(iban|kontonummer|bankverbindung|bic)\b'
    if re.search(bankdata_pattern, search_text, re.IGNORECASE):
        # Führe LLM-basierte Kontext-Analyse durch
        bankdata_check = check_bankdata_context_with_llm(text)

        # Nur wenn es sich um PRIVATE Bankdaten handelt, als sensibel markieren
        if bankdata_check['contains_private_bankdata']:
            result['contains_sensitive_data'] = True
            if 'BANKDATEN_PRIVAT' not in result['data_categories']:
                result['data_categories'].append('BANKDATEN_PRIVAT')

            # Füge DSGVO-Klassifizierung hinzu
            dsgvo_cat = "Art. 6 Abs. 1 DSGVO - Finanzdaten (natürliche Person)"
            if dsgvo_cat not in result['dsgvo_classification']:
                result['dsgvo_classification'].append(dsgvo_cat)

            # Setze Schutzklasse auf 'hoch'
            if highest_protection is None or protection_levels.get('hoch', 0) > protection_levels.get(highest_protection, 0):
                highest_protection = 'hoch'
                result['protection_level'] = highest_protection

            # Füge Details zur Begründung hinzu
            result['matched_keywords']['BANKDATEN_PRIVAT'] = [
                f"LLM-Analyse: {bankdata_check['context']} (Konfidenz: {bankdata_check['confidence']})"
            ]

    return result

def get_prompt_for_filetype(file_ext, summary_max_chars=1500):
    """
    Gibt einen RAG-optimierten, dateityp-spezifischen Prompt zurück.
    Optimiert für semantische Suche und Wissensextraktion.

    Args:
        file_ext: Dateierweiterung
        summary_max_chars: Maximale Länge der Zusammenfassung in Zeichen
    """
    # Basis-Prompt für RAG-Optimierung
    base_prompt = f"""Du bist ein System zur Wissensextraktion für semantische Suche (RAG).

Fasse den folgenden Dateiinhalt so zusammen, dass er für spätere Fragen maximal gut auffindbar und nutzbar ist.

REGELN:
- Maximal {summary_max_chars} Zeichen
- Sachlich, präzise, ohne Floskeln
- Keine Meta-Kommentare (z. B. „Diese Datei beschreibt…", „Zusammenfassung:", „Das Dokument enthält…")
- Keine Markdown-Formatierung (**, ##, -, etc.)
- Nur reiner Fließtext ohne Überschriften oder Listen
- Nutze klare, informationsdichte Sätze
- Behalte wichtige Fachbegriffe, Zahlen, Technologien und Personennamen
- Beschreibe Zweck, Inhalt, Kontext und Besonderheiten
- Falls vorhanden: Ziel, Funktion, Datenarten, Methoden, Abhängigkeiten

STRUKTUR (fließender Text ohne Überschriften):
- Worum geht es?
- Wozu dient es?
- Welche Inhalte/Daten/Logik sind enthalten?
- Was macht es besonders oder relevant?

PFLICHTFELD - KEYWORDS (auf neuer Zeile am Ende):
Die letzte Zeile MUSS folgendes Format haben:
Schlüsselbegriffe: Begriff1, Begriff2, Begriff3, Begriff4, Begriff5

Mindestens 3-8 zentrale Fachbegriffe, Technologien oder Themen als kommagetrennte Liste.
Die Keyword-Zeile MUSS mit "Schlüsselbegriffe:" beginnen.

WICHTIG: Antworte AUF DEUTSCH. Beginne direkt mit dem Inhalt, ohne Einleitung."""

    # Dateityp-spezifische Ergänzungen
    type_specific = {
        ".pdf": "Fokus: Dokumenteninhalt, Kernaussagen, Personen und ihre Rollen.",
        ".docx": "Fokus: Dokumenteninhalt, Kernaussagen, Personen und ihre Rollen.",
        ".doc": "Fokus: Dokumenteninhalt, Kernaussagen, Personen und ihre Rollen.",
        ".pptx": "Fokus: Präsentationsthemen, Kernbotschaften, Struktur der Folien.",
        ".ppt": "Fokus: Präsentationsthemen, Kernbotschaften, Struktur der Folien.",
        ".xlsx": "Fokus: Datenarten, Kategorien, Zweck der Tabelle, enthaltene Zahlen.",
        ".xls": "Fokus: Datenarten, Kategorien, Zweck der Tabelle, enthaltene Zahlen.",
        ".xlsm": "Fokus: Datenarten, Kategorien, Makro-Funktionalität, Automatisierung.",
        ".xltx": "Fokus: Vorlagenzweck, Struktur, verwendete Kategorien.",
        ".txt": "Fokus: Textinhalt, Zweck, enthaltene Informationen.",
        ".md": "Fokus: Dokumentstruktur, Hauptthemen, technische Details.",
        ".png": "Fokus: Bildinhalte, sichtbarer Text, Diagramme, Personen, Zweck.",
        ".jpg": "Fokus: Bildinhalte, sichtbare Personen, Kontext, Details.",
        ".jpeg": "Fokus: Bildinhalte, sichtbare Personen, Kontext, Details."
    }

    # Kombiniere Basis-Prompt mit dateityp-spezifischer Ergänzung
    specific = type_specific.get(file_ext, "Fokus: Inhalt, Zweck, Relevanz.")
    return f"{base_prompt}\n\n{specific}"

def extract_entities_with_lmstudio(text, file_path=None, file_ext=None):
    """
    Extrahiert Named Entities (Firmen, Personen, Institutionen, Organisationen) aus Text.
    Funktioniert sowohl für kurze als auch lange Texte.

    Args:
        text: Der zu analysierende Text
        file_path: Optional - Pfad zur Datei (für Bilder)
        file_ext: Optional - Dateierweiterung

    Returns:
        dict mit Listen: {'companies': [], 'persons': [], 'institutions': [], 'organizations': []}
    """
    # Prüfe ob es sich um eine Bilddatei handelt
    is_image = file_ext and file_ext.lower() in {".png", ".jpg", ".jpeg"}

    # Für Bilder mit Vision API
    if is_image and file_path:
        return extract_entities_from_image(file_path, file_ext)

    # Begrenze Text auf sinnvolle Länge für Entity-Extraktion
    # Für sehr lange Texte: verwende Anfang und Ende
    max_chars = 8000
    if len(text) > max_chars:
        # Nehme erste 6000 und letzte 2000 Zeichen
        truncated_text = text[:6000] + "\n...\n" + text[-2000:]
    else:
        truncated_text = text

    entity_prompt = """Extrahiere alle Named Entities aus dem folgenden Text.

KATEGORIEN:
- Firmen/Unternehmen: Namen von Firmen, Gesellschaften, Unternehmen
- Personen: Vollständige Namen von Personen (Vor- und Nachname wenn möglich)
- Institutionen: Behörden, Ämter, staatliche Einrichtungen, Bildungseinrichtungen
- Organisationen: Vereine, Verbände, NGOs, andere Organisationen

REGELN:
- Extrahiere nur tatsächlich im Text vorkommende Namen
- Keine generischen Begriffe wie "der Kunde", "das Unternehmen"
- Vollständige Namen bevorzugen
- Keine Duplikate
- Falls keine Entitäten in einer Kategorie: leere Liste

AUSGABEFORMAT (exakt so):
FIRMEN: Firma1, Firma2, Firma3
PERSONEN: Max Mustermann, Erika Beispiel
INSTITUTIONEN: Bundesamt für XY, Universität Z
ORGANISATIONEN: Verein ABC, Verband DEF

WICHTIG:
- Wenn eine Kategorie leer ist, schreibe: "FIRMEN:" (ohne Einträge)
- Trenne mehrere Einträge mit Komma
- Antworte AUF DEUTSCH
- Verwende exakt das Format oben"""

    payload = {
        "model": MODEL_NAME,
        "messages": [
            {
                "role": "system",
                "content": "Du bist ein System zur Extraktion von Named Entities. Extrahiere nur tatsächlich vorhandene Namen in den angegebenen Kategorien."
            },
            {
                "role": "user",
                "content": f"{entity_prompt}\n\nTEXT:\n{truncated_text}"
            }
        ],
        "temperature": 0.1,  # Niedrige Temperatur für konsistente Extraktion
        "max_tokens": 500,
    }

    try:
        resp = requests.post(LMSTUDIO_API_URL, json=payload, timeout=120)
        resp.raise_for_status()

        data = resp.json()
        response_text = data["choices"][0]["message"]["content"]

        # Parse die strukturierte Antwort
        entities = parse_entity_response(response_text)
        return entities

    except Exception as e:
        print(f"  → Warnung: Entity-Extraktion fehlgeschlagen: {str(e)[:100]}")
        # Fallback: Leere Listen
        return {
            'companies': [],
            'persons': [],
            'institutions': [],
            'organizations': []
        }

def extract_entities_from_image(image_path, file_ext):
    """
    Extrahiert Named Entities aus einem Bild mit Vision API.

    Returns:
        dict mit Listen: {'companies': [], 'persons': [], 'institutions': [], 'organizations': []}
    """
    import base64

    try:
        # Lese Bild und konvertiere zu Base64
        with open(image_path, 'rb') as img_file:
            img_data = base64.b64encode(img_file.read()).decode('utf-8')

        # Bestimme MIME-Type
        mime_type = "image/png" if file_ext.lower() == ".png" else "image/jpeg"

        entity_prompt = """Extrahiere alle sichtbaren Named Entities aus diesem Bild.

KATEGORIEN:
- Firmen/Unternehmen: Namen von Firmen, Gesellschaften, Unternehmen, Logos
- Personen: Namen von Personen (wenn lesbar/erkennbar)
- Institutionen: Behörden, Ämter, staatliche Einrichtungen, Bildungseinrichtungen
- Organisationen: Vereine, Verbände, NGOs, andere Organisationen

AUSGABEFORMAT (exakt so):
FIRMEN: Firma1, Firma2
PERSONEN: Max Mustermann, Erika Beispiel
INSTITUTIONEN: Bundesamt für XY
ORGANISATIONEN: Verein ABC

Falls eine Kategorie keine Einträge hat, lasse sie leer (z.B. "FIRMEN:")"""

        payload = {
            "model": MODEL_NAME,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": entity_prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{img_data}"
                            }
                        }
                    ]
                }
            ],
            "temperature": 0.1,
            "max_tokens": 300,
        }

        resp = requests.post(LMSTUDIO_API_URL, json=payload, timeout=180)
        resp.raise_for_status()

        data = resp.json()
        response_text = data["choices"][0]["message"]["content"]

        # Parse die strukturierte Antwort
        entities = parse_entity_response(response_text)
        return entities

    except Exception as e:
        print(f"  → Warnung: Entity-Extraktion aus Bild fehlgeschlagen: {str(e)[:100]}")
        return {
            'companies': [],
            'persons': [],
            'institutions': [],
            'organizations': []
        }

def parse_entity_response(response_text):
    """
    Parst die strukturierte Entity-Antwort vom LLM.

    Erwartet Format:
    FIRMEN: Firma1, Firma2
    PERSONEN: Person1, Person2
    INSTITUTIONEN: Institution1
    ORGANISATIONEN: Org1, Org2

    Returns:
        dict: {'companies': [...], 'persons': [...], 'institutions': [...], 'organizations': [...]}
    """
    import re

    entities = {
        'companies': [],
        'persons': [],
        'institutions': [],
        'organizations': []
    }

    # Mapping von deutschen Labels zu dict keys
    label_mapping = {
        'FIRMEN': 'companies',
        'PERSONEN': 'persons',
        'INSTITUTIONEN': 'institutions',
        'ORGANISATIONEN': 'organizations',
        # Auch englische Varianten für Robustheit
        'COMPANIES': 'companies',
        'PERSONS': 'persons',
        'INSTITUTIONS': 'institutions',
        'ORGANIZATIONS': 'organizations',
        # Weitere mögliche Varianten
        'UNTERNEHMEN': 'companies',
        'FIRMA': 'companies',
    }

    # Parse Zeile für Zeile
    lines = response_text.strip().split('\n')

    for line in lines:
        line = line.strip()
        if not line or ':' not in line:
            continue

        # Trenne Label und Inhalt
        parts = line.split(':', 1)
        if len(parts) != 2:
            continue

        label = parts[0].strip().upper()
        content = parts[1].strip()

        # Finde passendes Mapping
        entity_key = label_mapping.get(label)
        if not entity_key:
            continue

        # Parse kommagetrennte Einträge
        if content:
            items = [item.strip() for item in content.split(',') if item.strip()]
            # Entferne Duplikate und leere Einträge
            items = list(dict.fromkeys(items))  # Erhält Reihenfolge und entfernt Duplikate
            entities[entity_key].extend(items)

    # Entferne finale Duplikate über alle geparsten Zeilen hinweg
    for key in entities:
        entities[key] = list(dict.fromkeys(entities[key]))

    return entities

def summarize_image_with_lmstudio(image_path, file_ext):
    """Analysiert ein Bild mit der Vision API von LM Studio."""
    import base64

    try:
        # Lese Bild und konvertiere zu Base64
        with open(image_path, 'rb') as img_file:
            img_data = base64.b64encode(img_file.read()).decode('utf-8')

        # Bestimme MIME-Type
        mime_type = "image/png" if file_ext.lower() == ".png" else "image/jpeg"

        # Hole dateityp-spezifischen Prompt
        user_prompt = get_prompt_for_filetype(file_ext)

        payload = {
            "model": MODEL_NAME,
            "messages": [
                {
                    "role": "system",
                    "content": "Du bist ein Wissensextraktionssystem für semantische Suche. Erstelle informationsdichte Zusammenfassungen in reinem Fließtext ohne Meta-Kommentare, ohne Markdown-Formatierung und ohne Überschriften. Fokussiere auf Fakten, Zahlen, Namen und Fachbegriffe. Beginne direkt mit dem Inhalt. WICHTIG: Gib KEINE Gedankenprozesse oder [THINK]-Tags aus, nur die finale Zusammenfassung."
                },
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": user_prompt},
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{img_data}"
                            }
                        }
                    ]
                }
            ],
            "temperature": 0.3,
            "max_tokens": 400,  # Erhöht für ~1000 Zeichen Output
        }

        resp = requests.post(LMSTUDIO_API_URL, json=payload, timeout=300)
        resp.raise_for_status()

        data = resp.json()
        summary = data["choices"][0]["message"]["content"]

        # Entferne [THINK] Tags von Reasoning-Modellen
        # Reasoning-Modelle wie ministral-3-14b-reasoning umschließen ihre Gedanken mit [THINK]...[/THINK]
        import re
        # Entferne alles zwischen [THINK] und [/THINK]
        summary = re.sub(r'\[THINK\].*?\[/THINK\]', '', summary, flags=re.DOTALL)
        # Entferne verbleibende einzelne Tags
        summary = summary.replace('[THINK]', '').replace('[/THINK]', '')
        # Bereinige mehrfache Leerzeichen und Zeilenumbrüche
        summary = re.sub(r'\n{3,}', '\n\n', summary)  # Max 2 Zeilenumbrüche
        summary = summary.strip()

        return summary

    except Exception as e:
        print(f"Fehler bei Bildanalyse: {e}")
        # Fallback: Gebe einen Platzhalter zurück
        return f"Bilddatei ({file_ext}). Vision-Analyse fehlgeschlagen: {str(e)[:100]}"

# Globaler Lern-Cache für erfolgreiche Context-Größen
# Strategie: Graduelle Aufwärts-Exploration mit adaptivem Lernen
# Struktur: {model_name: {'current_max': int, 'successes': [int], 'consecutive_ok': int, 'last_failed': int}}
_LEARNED_MAX_CHARS = {}

def summarize_with_lmstudio(text, file_path=None, file_ext=None, max_chars=30000, summary_max_chars=1500):
    # Adaptive Textkürzung mit automatischem Retry bei Context-Overflow
    # ministral-3-14b-reasoning hat größeres Context-Fenster
    # Start mit ~30000 Zeichen (~7500 Tokens), bei Fehler schrittweise reduzieren

    # Stelle sicher, dass file_ext ein String ist
    if file_ext and not isinstance(file_ext, str):
        raise TypeError(f"file_ext muss ein String sein, nicht {type(file_ext)}")

    # Prüfe ob es sich um eine Bilddatei handelt
    is_image = file_ext and file_ext.lower() in {".png", ".jpg", ".jpeg"}

    if is_image and file_path:
        # Für Bilder: Verwende Vision API
        return summarize_image_with_lmstudio(file_path, file_ext)

    # Entferne problematische Zeichen und normalisiere Whitespace
    text = text.strip()
    if not text:
        raise ValueError("Text ist leer nach Bereinigung")

    # Wenn der Text kürzer als die Zielgröße ist, kopiere ihn direkt
    if len(text) <= summary_max_chars:
        print(f"Text ({len(text)} Zeichen) ist kürzer als Zielgröße ({summary_max_chars}), kopiere Original")
        return text

    # Versuche mit verschiedenen Textlängen, falls Context zu groß ist
    # Berechne retry_lengths basierend auf MAX_CONTEXT_TOKENS
    # Annahme: ~4 Zeichen pro Token (konservativ für deutsche Texte)
    chars_per_token = 4
    max_chars = (MAX_CONTEXT_TOKENS - 1000) * chars_per_token  # Reserve 1000 Tokens für Prompt und Antwort

    actual_text_length = len(text)

    # Adaptive Lernlogik mit gradueller Aufwärts-Exploration
    global _LEARNED_MAX_CHARS

    # Initialisiere Lern-Struktur für dieses Modell
    if MODEL_NAME not in _LEARNED_MAX_CHARS:
        _LEARNED_MAX_CHARS[MODEL_NAME] = {
            'current_max': max_chars // 2,  # Start konservativ bei 50% vom Maximum
            'successes': [],                # Liste der letzten 10 erfolgreichen Größen
            'consecutive_ok': 0,            # Zähler für aufeinanderfolgende Erfolge
            'last_failed': None             # Letzte fehlgeschlagene Größe (obere Grenze)
        }

    learned_data = _LEARNED_MAX_CHARS[MODEL_NAME]

    # Berechne Startpunkt basierend auf Lernhistorie
    # Wenn wir mehrere Erfolge hatten, versuche schrittweise nach oben zu gehen
    if learned_data['consecutive_ok'] >= 3:
        # Nach 3 aufeinanderfolgenden Erfolgen: Erhöhe um 10%
        exploration_max = int(learned_data['current_max'] * 1.10)
        # Aber nicht über bekannte Fehlergrenze hinaus
        if learned_data['last_failed']:
            exploration_max = min(exploration_max, int(learned_data['last_failed'] * 0.95))
        else:
            exploration_max = min(exploration_max, max_chars)

        print(f"  → 🔼 Exploration: Teste größeren Context ({exploration_max:,} Zeichen, +10%)")
        start_chars = exploration_max
    else:
        # Nutze aktuell bekannte sichere Größe
        start_chars = learned_data['current_max']
        if learned_data['successes']:
            avg_success = int(sum(learned_data['successes']) / len(learned_data['successes']))
            print(f"  → Nutze gelernte Größe: {start_chars:,} Zeichen (Ø {avg_success:,})")

    # Erstelle Retry-Liste: Start mit optimistischer Größe, dann sanfte Reduktion
    retry_lengths = []

    # Primärversuch: Exploration oder gelernte Größe
    base_chars = min(start_chars, actual_text_length)
    retry_lengths.append(base_chars)

    # Fallback-Schritte bei Fehler: -15%, -30%, -45%, -60%, dann drastischer
    # Sanftere Schritte als vorher für bessere Konvergenz
    fallback_steps = [0.85, 0.70, 0.55, 0.40, 0.30, 0.20]
    for step in fallback_steps:
        chars = int(base_chars * step)
        if chars > summary_max_chars and chars not in retry_lengths:
            retry_lengths.append(chars)

    # Absolutes Minimum als letzte Rettung
    min_fallback = min(3000, actual_text_length)
    if min_fallback not in retry_lengths:
        retry_lengths.append(min_fallback)

    # Sortiere absteigend
    retry_lengths = sorted(list(set(retry_lengths)), reverse=True)

    # Hole dateityp-spezifischen Prompt
    user_prompt = get_prompt_for_filetype(file_ext, summary_max_chars) if file_ext else get_prompt_for_filetype("", summary_max_chars)

    # Berechne max_tokens basierend auf Zielgröße
    # ~2.5 Zeichen pro Token für deutsche Texte
    max_tokens = int(summary_max_chars / 2.5) + 50  # +50 für Keywords

    for attempt, current_max_chars in enumerate(retry_lengths, 1):
        truncated_text = text[:current_max_chars]

        payload = {
            "model": MODEL_NAME,
            "messages": [
                {
                    "role": "system",
                    "content": "Du bist ein Wissensextraktionssystem für semantische Suche. Erstelle informationsdichte Zusammenfassungen in reinem Fließtext ohne Meta-Kommentare (z.B. 'Zusammenfassung:', 'Diese Datei...'), ohne Markdown-Formatierung (**, ##, -) und ohne Überschriften. Fokussiere auf Fakten, Zahlen, Namen und Fachbegriffe. Beginne direkt mit dem Inhalt."
                },
                {
                    "role": "user",
                    "content": f"{user_prompt}\n\nDokument:\n{truncated_text}"
                },
            ],
            "temperature": 0.3,
            "max_tokens": max_tokens,  # Dynamisch basierend auf SUMMARY_MAX_CHARS
        }

        try:
            resp = requests.post(LMSTUDIO_API_URL, json=payload, timeout=300)
            resp.raise_for_status()

            # Erfolg! Gib die Zusammenfassung zurück
            data = resp.json()
            summary = data["choices"][0]["message"]["content"]

            # Adaptive Lernlogik: Aktualisiere basierend auf Erfolg
            learned_data = _LEARNED_MAX_CHARS[MODEL_NAME]

            if attempt == 1:
                # Erfolg beim ersten Versuch
                learned_data['consecutive_ok'] += 1

                # Bei Exploration: Aktualisiere current_max nach oben
                if current_max_chars > learned_data['current_max']:
                    old_max = learned_data['current_max']
                    learned_data['current_max'] = current_max_chars
                    print(f"  ✓ Erfolg bei Exploration! Neue Grenze: {current_max_chars:,} (vorher: {old_max:,})")
                else:
                    print(f"  ✓ Erfolg mit {current_max_chars:,} Zeichen")

            else:
                # Erfolg nach Retry - Reset consecutive counter
                learned_data['consecutive_ok'] = 0
                # Erfolg bei niedrigerer Größe: Aktualisiere current_max konservativ
                learned_data['current_max'] = current_max_chars
                print(f"  ✓ Erfolgreich mit {current_max_chars:,} Zeichen (Versuch {attempt})")

            # Speichere in Success-Historie (behalte nur letzte 10)
            learned_data['successes'].append(current_max_chars)
            if len(learned_data['successes']) > 10:
                learned_data['successes'].pop(0)

            return summary

        except requests.exceptions.HTTPError as e:
            # Prüfe ob es ein Context-Overflow-Fehler ist
            try:
                error_data = resp.json()
                error_msg = str(error_data.get("error", ""))

                # Context-Overflow erkannt? Prüfe auf verschiedene Fehlermeldungen
                is_context_error = (
                    "context" in error_msg.lower() or
                    "token" in error_msg.lower() or
                    "length" in error_msg.lower()
                )

                if is_context_error:
                    # Lernlogik: Speichere Fehlergrenze
                    learned_data = _LEARNED_MAX_CHARS[MODEL_NAME]

                    # Merke diese Größe als "zu groß"
                    if learned_data['last_failed'] is None or current_max_chars < learned_data['last_failed']:
                        learned_data['last_failed'] = current_max_chars

                    # Reset consecutive successes
                    learned_data['consecutive_ok'] = 0

                    if attempt < len(retry_lengths):
                        # Berechne geschätzte Tokens für Debug-Ausgabe
                        estimated_tokens = current_max_chars // 4
                        print(f"  ✗ Context-Limit erreicht ({current_max_chars:,} Zeichen ≈ {estimated_tokens:,} Tokens)")
                        print(f"     Grenze gespeichert, versuche mit weniger...")
                        continue  # Nächster Versuch mit weniger Text
                    else:
                        print(f"  → Alle Retry-Versuche fehlgeschlagen")
                        raise ValueError(f"Text zu lang selbst nach {len(retry_lengths)} Versuchen: {error_msg}")
                else:
                    # Anderer HTTP-Fehler
                    print(f"HTTP-Fehler {resp.status_code}:")
                    print(f"Response-Text: {resp.text}")
                    raise
            except (ValueError, KeyError, json.JSONDecodeError):
                # Kein JSON oder kein error-Feld - könnte trotzdem Context-Fehler sein
                if resp.status_code == 400 and attempt < len(retry_lengths):
                    estimated_tokens = current_max_chars // 4
                    print(f"  → HTTP 400 Fehler ({current_max_chars:,} Zeichen ≈ {estimated_tokens:,} Tokens), versuche mit weniger...")
                    print(f"     Response: {resp.text[:150]}...")  # Erste 150 Zeichen der Response
                    continue
                else:
                    print(f"HTTP-Fehler {resp.status_code}:")
                    print(f"Response-Text: {resp.text}")
                    raise

    # Falls alle Versuche fehlschlagen
    raise ValueError("Zusammenfassung fehlgeschlagen nach allen Retry-Versuchen")

def process_file(src_file):
    """
    Verarbeitet eine einzelne Datei und erstellt eine JSON-Zusammenfassung.

    Returns:
        dict: OCR-Informationen falls verfügbar, sonst None
    """
    rel_path = os.path.relpath(src_file, SRC_ROOT)
    dst_dir = os.path.join(DST_ROOT, os.path.dirname(rel_path))
    os.makedirs(dst_dir, exist_ok=True)

    # Summary-Datei neben die Quelle legen, aber unter D:\LLM
    dst_file = os.path.join(dst_dir, os.path.basename(src_file) + ".json")

    # Prüfe ob Datei existiert und valide ist
    if os.path.exists(dst_file):
        if validate_json_file(dst_file, src_file):
            print("Überspringe (valide Summary existiert):", dst_file)
            # Lese OCR-Info aus existierender JSON-Datei für Statistik
            try:
                with open(dst_file, 'r', encoding='utf-8') as f:
                    existing_data = json.load(f)
                    return existing_data.get('ocr_info', None)
            except:
                return None
        else:
            print("Lösche fehlerhafte oder veraltete JSON-Datei:", dst_file)
            try:
                os.remove(dst_file)
            except Exception as e:
                print(f"Fehler beim Löschen von {dst_file}: {e}")
                return None

    path_obj = pathlib.Path(src_file)
    print("Verarbeite:", src_file)

    # Prüfe ob Datei zugänglich ist
    if not is_file_accessible(src_file):
        print(f"Überspringe Datei, da nicht zugänglich: {src_file}")
        return None

    # Für Bilddateien: Prüfe Mindestgröße (ignoriere kleine Icons)
    file_ext = path_obj.suffix.lower()
    if file_ext in {".png", ".jpg", ".jpeg"}:
        try:
            file_size = os.path.getsize(src_file)
            if file_size < MIN_IMAGE_SIZE:
                print(f"Überspringe kleine Bilddatei ({file_size} Bytes < {MIN_IMAGE_SIZE} Bytes): {src_file}")
                return None
        except OSError as e:
            print(f"Fehler beim Prüfen der Dateigröße: {e}")
            return None

    try:
        result = extract_text(path_obj)
        # Stelle sicher, dass wir ein Tuple bekommen
        if isinstance(result, tuple) and len(result) == 2:
            text, ocr_info = result
        else:
            # Fallback für unerwartetes Format
            print(f"Unerwartetes Format von extract_text(): {type(result)}")
            if isinstance(result, str):
                text = result
                ocr_info = None
            else:
                print(f"Fehler: Kann Text nicht extrahieren, unbekanntes Format: {result}")
                return None
    except Exception as e:
        print(f"Fehler beim Extrahieren von Text aus {src_file}: {e}")
        import traceback
        traceback.print_exc()
        return None

    # file_ext wurde bereits oben definiert (Zeile 425)
    is_image = file_ext in {".png", ".jpg", ".jpeg"}

    # Stelle sicher, dass text ein String ist
    if not isinstance(text, str):
        print(f"Fehler: Text ist kein String sondern {type(text)}: {text}")
        return None

    if not is_image and not text.strip():
        # Prüfe ob das Problem fehlende OCR-Unterstützung ist
        if ocr_info and not ocr_info.get('used_ocr') and not OCR_AVAILABLE:
            # Dies ist wahrscheinlich eine gescannte PDF ohne verfügbares OCR
            print("!" * 70)
            print("ÜBERSPRUNGEN: Gescannte PDF ohne OCR-Unterstützung")
            print("!" * 70)
            print(f"Datei: {src_file}")
            print("\nDiese Datei scheint gescannten Text zu enthalten und benötigt OCR.")
            print("OCR ist nicht verfügbar (pytesseract/Tesseract nicht installiert).")
            print("\nInstallation:")
            print("  macOS:  brew install tesseract tesseract-lang")
            print("  Linux:  sudo apt-get install tesseract-ocr tesseract-ocr-deu")
            print("  Python: pip install pytesseract pillow")
            print("!" * 70)
        else:
            print("Kein Text extrahiert, überspringe:", src_file)
        return None

    if not is_image:
        print(f"Text extrahiert: {len(text)} Zeichen")

    try:
        # Debug: Prüfe file_ext Typ
        if not isinstance(file_ext, str):
            print(f"FEHLER: file_ext hat falschen Typ: {type(file_ext)}, Wert: {file_ext}")
            return None

        # Übergebe file_path und file_ext für dateityp-spezifische Verarbeitung
        summary = summarize_with_lmstudio(text, file_path=src_file, file_ext=file_ext, summary_max_chars=SUMMARY_MAX_CHARS)

        # Zeige die ersten 100 Zeichen der Zusammenfassung
        summary_preview = summary[:100] + "..." if len(summary) > 100 else summary
        print(f"Zusammenfassung: {summary_preview}")
    except ValueError as e:
        error_msg = f"Validierungsfehler: {e}"
        action = ask_on_lmstudio_error(error_msg, src_file)
        if action == "abort":
            raise SystemExit("Verarbeitung durch Benutzer abgebrochen.")
        return None
    except TypeError as e:
        error_msg = f"Typfehler: {e}"
        action = ask_on_lmstudio_error(error_msg, src_file)
        if action == "abort":
            raise SystemExit("Verarbeitung durch Benutzer abgebrochen.")
        return None
    except requests.exceptions.RequestException as e:
        error_msg = f"Netzwerkfehler bei der Zusammenfassung: {e}"
        action = ask_on_lmstudio_error(error_msg, src_file)
        if action == "abort":
            raise SystemExit("Verarbeitung durch Benutzer abgebrochen.")
        return None

    # Sammle Datei-Metadaten
    stat = os.stat(src_file)

    # Extrahiere Named Entities aus dem Text
    # Dies geschieht für ALLE Texte, egal ob kurz oder lang
    print("Extrahiere Named Entities...")
    entities = extract_entities_with_lmstudio(text, file_path=src_file, file_ext=file_ext)

    # Extrahiere zusätzliche Entities aus dem Dateipfad
    path_entities = extract_entities_from_path(src_file)

    # Merge Pfad-Entities mit Text-Entities (ohne Duplikate)
    for company in path_entities['companies']:
        if company not in entities['companies']:
            entities['companies'].append(company)

    # Speichere Projektnamen separat (neues Feld)
    entities['projects'] = path_entities['projects']

    # Extrahiere Kontaktinformationen (URLs, E-Mails, Telefon) - Regex-basiert, sehr schnell
    print("Extrahiere Kontaktinformationen...")
    contact_info = extract_contact_info_from_text(text)
    entities['urls'] = contact_info['urls']
    entities['emails'] = contact_info['emails']
    entities['phone_numbers'] = contact_info['phone_numbers']

    # Zeige gefundene Entities (wenn vorhanden)
    entity_count = sum(len(v) for v in entities.values())
    if entity_count > 0:
        print(f"  → Gefunden: {len(entities['companies'])} Firmen, {len(entities['persons'])} Personen, "
              f"{len(entities['institutions'])} Institutionen, {len(entities['organizations'])} Organisationen, "
              f"{len(entities.get('projects', []))} Projekte, {len(entities.get('urls', []))} URLs, "
              f"{len(entities.get('emails', []))} E-Mails, {len(entities.get('phone_numbers', []))} Telefonnummern")

    # Extrahiere Schlüsselbegriffe aus der Zusammenfassung
    # Die Schlüsselbegriffe sollten am Ende der Zusammenfassung stehen
    keywords = []
    summary_text = summary

    # Suche nach Keyword-Markern wie "Schlüsselbegriffe:", "Keywords:", etc.
    import re

    # Muster für verschiedene Keyword-Marker (auch mit Absatz/Newline davor)
    keyword_patterns = [
        r'\n\s*Schlüsselbegriffe:\s*(.+?)$',
        r'\n\s*Keywords?:\s*(.+?)$',
        r'\n\s*Zentrale Begriffe:\s*(.+?)$',
        # Fallback: Suche auch ohne Newline am Anfang
        r'Schlüsselbegriffe:\s*(.+?)$',
        r'Keywords?:\s*(.+?)$',
    ]

    for pattern in keyword_patterns:
        match = re.search(pattern, summary, re.IGNORECASE | re.MULTILINE)
        if match:
            keyword_string = match.group(1).strip()
            # Extrahiere kommagetrennte Keywords
            if ',' in keyword_string:
                keywords = [kw.strip() for kw in keyword_string.split(',') if kw.strip()]
                # Entferne die Keyword-Zeile aus der Zusammenfassung
                summary_text = re.sub(pattern, '', summary, flags=re.IGNORECASE | re.MULTILINE).strip()
                break

    # Fallback: Wenn keine Keywords gefunden wurden, versuche letzte Zeile
    if not keywords:
        lines = summary.strip().split('\n')
        if len(lines) > 1:
            # Letzte Zeile könnte die Keywords enthalten
            last_line = lines[-1].strip()
            # Prüfe ob die letzte Zeile hauptsächlich aus kommagetrennten Wörtern besteht
            if ',' in last_line and len(last_line) < 300:
                # Extrahiere Keywords
                keywords = [kw.strip() for kw in last_line.split(',') if kw.strip()]
                # Entferne die Keyword-Zeile aus der Zusammenfassung
                summary_text = '\n'.join(lines[:-1]).strip()

    # Berechne Content-Hash für Änderungserkennung
    content_hash = calculate_content_hash(src_file)

    # Klassifiziere sensible/schutzbedürftige Daten gemäß DSGVO/BDSG
    print("Klassifiziere DSGVO-relevante Inhalte...")
    sensitive_classification = classify_sensitive_data(text, file_path=src_file)

    # Zeige Klassifizierungsergebnis
    if sensitive_classification['contains_sensitive_data']:
        print(f"  ⚠️  DSGVO-WARNUNG: Besonders schutzbedürftige Daten erkannt!")
        print(f"      Kategorien: {', '.join(sensitive_classification['data_categories'])}")
        print(f"      Schutzklasse: {sensitive_classification['protection_level']}")

    metadata = {
        "path": rel_path,
        "ext": path_obj.suffix.lower(),
        "size": stat.st_size,
        "created": datetime.fromtimestamp(stat.st_ctime).isoformat(),
        "modified": datetime.fromtimestamp(stat.st_mtime).isoformat(),
        "content_hash": content_hash,
        "chars": len(text),
        "summary": summary_text,
        "keywords": keywords,
        "entities": {
            "companies": entities['companies'],
            "persons": entities['persons'],
            "institutions": entities['institutions'],
            "organizations": entities['organizations'],
            "projects": entities.get('projects', []),
            "urls": entities.get('urls', []),
            "emails": entities.get('emails', []),
            "phone_numbers": entities.get('phone_numbers', [])
        },
        # DSGVO/BDSG-Klassifizierung
        "dsgvo_classification": {
            "contains_sensitive_data": sensitive_classification['contains_sensitive_data'],
            "data_categories": sensitive_classification['data_categories'],
            "legal_basis": sensitive_classification['dsgvo_classification'],
            "protection_level": sensitive_classification['protection_level'],
            "detected_keywords": sensitive_classification['matched_keywords']
        }
    }

    # Füge OCR-Info hinzu falls verfügbar
    if ocr_info and ocr_info.get('used_ocr'):
        metadata['ocr_info'] = ocr_info

    with open(dst_file, "w", encoding="utf-8") as f:
        json.dump(metadata, f, ensure_ascii=False, indent=2)

    print(f"Summary erfolgreich erstellt: {dst_file}")

    return ocr_info

def validate_phone_number(phone):
    """
    Validiert eine Telefonnummer gegen die aktuellen strengen Regex-Pattern.
    Filtert falsche Positives wie Projektnummern (091-2024) heraus.

    Args:
        phone: Telefonnummer als String

    Returns:
        True wenn gültig, False wenn ungültig
    """
    import re

    # Normalisiere für Prüfung
    digits_only = re.sub(r'\D', '', phone)

    # Prüfe Mindestlänge (8 Ziffern)
    if len(digits_only) < 8:
        return False

    # Muss mit 0 oder 49 beginnen (deutsche Nummern)
    if not (digits_only.startswith('0') or digits_only.startswith('49')):
        return False

    # Prüfe gegen Pattern (mindestens eins muss matchen)
    phone_patterns = [
        r'\+49[\s\-]?\(?\d{2,4}\)?[\s\-]?\d{3,10}',
        r'\+49[\s\-]?\d{2,4}[\s\-/]\d{6,10}',
        r'0049[\s\-]?\d{2,4}[\s\-]?\d{6,10}',
        r'\(0\d{2,4}\)[\s\-]?\d{6,10}',
        r'0\d{2,4}[\s\-/]\d{6,10}',
        r'0\d{9,11}',
    ]

    for pattern in phone_patterns:
        # Nutze search statt fullmatch für flexiblere Matching
        if re.search(f'^{pattern}$', phone.strip()):
            return True

    return False

def update_json_with_contact_info(json_path, src_file_path):
    """
    Trägt fehlende Kontaktinformationen (URLs, E-Mails, Telefon) in existierender JSON nach.
    Validiert auch vorhandene Telefonnummern und entfernt ungültige (z.B. Projektnummern).
    Sehr schnell (nur Regex, kein LLM), spart massive Prozesszeit.

    Args:
        json_path: Pfad zur JSON-Datei
        src_file_path: Pfad zur Quelldatei (zum Text-Extrahieren)

    Returns:
        True wenn Update durchgeführt wurde, False wenn nicht nötig
    """
    try:
        # Lese JSON
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # Prüfe ob Kontaktfelder fehlen oder ungültige Telefonnummern enthalten
        entities = data.get('entities', {})

        # Prüfe fehlende Felder
        needs_extraction = (
            'urls' not in entities or
            'emails' not in entities or
            'phone_numbers' not in entities
        )

        # Prüfe vorhandene Telefonnummern auf Validität
        needs_phone_validation = False
        invalid_phones = []
        if 'phone_numbers' in entities and entities['phone_numbers']:
            for phone in entities['phone_numbers']:
                if not validate_phone_number(phone):
                    invalid_phones.append(phone)
                    needs_phone_validation = True

        # Wenn nichts zu tun ist, überspringe
        if not needs_extraction and not needs_phone_validation:
            return False

        # Extrahiere Text aus Quelldatei (verwende existierende Funktionen)
        path_obj = pathlib.Path(src_file_path)
        file_ext = path_obj.suffix.lower()

        # Nutze die bestehenden Extract-Funktionen
        if file_ext == ".pdf":
            text, _ = extract_text_pdf(src_file_path)
        elif file_ext in {".docx", ".doc"}:
            text = extract_text_docx(src_file_path)
        elif file_ext in {".pptx", ".ppt"}:
            text = extract_text_pptx(src_file_path)
        elif file_ext in {".xlsx", ".xls", ".xlsm", ".xltx"}:
            text = extract_text_xlsx(src_file_path)
        elif file_ext in {".txt", ".md"}:
            text = extract_text_txt(src_file_path)
        else:
            # Bilder oder unbekannt - überspringe
            return False

        # Extrahiere Kontaktinformationen (wenn nötig)
        if needs_extraction or needs_phone_validation:
            contact_info = extract_contact_info_from_text(text)

            # Aktualisiere fehlende Felder
            if 'urls' not in entities:
                entities['urls'] = contact_info['urls']
            if 'emails' not in entities:
                entities['emails'] = contact_info['emails']

            # Telefonnummern: Entweder nachtragen oder neu extrahieren (wenn ungültige gefunden)
            if 'phone_numbers' not in entities:
                entities['phone_numbers'] = contact_info['phone_numbers']
            elif needs_phone_validation:
                # Entferne ungültige und füge neu extrahierte hinzu
                valid_existing = [p for p in entities['phone_numbers'] if validate_phone_number(p)]
                # Kombiniere mit neu extrahierten (ohne Duplikate)
                all_phones = valid_existing + [p for p in contact_info['phone_numbers'] if p not in valid_existing]
                entities['phone_numbers'] = all_phones

                if invalid_phones:
                    print(f"  🧹 Entfernt {len(invalid_phones)} ungültige Telefonnummern: {invalid_phones[:3]}{'...' if len(invalid_phones) > 3 else ''}")

            # Speichere aktualisierte JSON
            with open(json_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)

            if needs_extraction:
                print(f"  ⚡ Kontaktinformationen nachgetragen: {len(contact_info['urls'])} URLs, "
                      f"{len(contact_info['emails'])} E-Mails, {len(contact_info['phone_numbers'])} Telefonnummern")
            elif needs_phone_validation:
                print(f"  ✓ Telefonnummern validiert: {len(entities['phone_numbers'])} gültig, {len(invalid_phones)} entfernt")

            return True

        return False

    except Exception as e:
        print(f"Fehler beim Nachtragen der Kontaktinformationen: {e}")
        return False

def update_json_with_dsgvo_classification(json_path, src_file_path):
    """
    Trägt DSGVO-Klassifizierung in existierende JSON-Dateien nach.
    Sehr schnell (nur Regex, kein LLM), analysiert Text auf sensible Daten.

    Args:
        json_path: Pfad zur JSON-Datei
        src_file_path: Pfad zur Quelldatei (zum Text-Extrahieren)

    Returns:
        True wenn Update durchgeführt wurde, False wenn nicht nötig
    """
    try:
        # Lese JSON
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # Prüfe ob DSGVO-Klassifizierung fehlt oder veraltet ist
        needs_classification = 'dsgvo_classification' not in data

        # Wenn nichts zu tun ist, überspringe
        if not needs_classification:
            return False

        # Extrahiere Text aus Quelldatei (verwende existierende Funktionen)
        path_obj = pathlib.Path(src_file_path)
        file_ext = path_obj.suffix.lower()

        # Nutze die bestehenden Extract-Funktionen
        text = ""
        try:
            if file_ext == ".pdf":
                text, _ = extract_text_pdf(src_file_path)
            elif file_ext in {".docx", ".doc"}:
                text = extract_text_docx(src_file_path)
            elif file_ext in {".pptx", ".ppt"}:
                text = extract_text_pptx(src_file_path)
            elif file_ext in {".xlsx", ".xls", ".xlsm", ".xltx"}:
                text = extract_text_xlsx(src_file_path)
            elif file_ext in {".txt", ".md"}:
                with open(src_file_path, 'r', encoding='utf-8', errors='ignore') as f:
                    text = f.read()
            elif file_ext in {".png", ".jpg", ".jpeg"}:
                # Für Bilder: Verwende Summary falls vorhanden
                text = data.get('summary', '')
            else:
                print(f"Unbekannter Dateityp für DSGVO-Update: {file_ext}")
                return False
        except Exception as e:
            print(f"Fehler beim Text-Extrahieren für DSGVO-Klassifizierung: {e}")
            return False

        if not text or not text.strip():
            # Kein Text verfügbar - verwende Summary als Fallback
            text = data.get('summary', '')

        if not text or not text.strip():
            print("  ⚠️  Kein Text für DSGVO-Klassifizierung verfügbar")
            return False

        # Klassifiziere
        sensitive_classification = classify_sensitive_data(text, file_path=src_file_path)

        # Füge DSGVO-Klassifizierung hinzu
        data['dsgvo_classification'] = {
            "contains_sensitive_data": sensitive_classification['contains_sensitive_data'],
            "data_categories": sensitive_classification['data_categories'],
            "legal_basis": sensitive_classification['dsgvo_classification'],
            "protection_level": sensitive_classification['protection_level'],
            "detected_keywords": sensitive_classification['matched_keywords']
        }

        # Speichere aktualisierte JSON
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

        # Zeige Ergebnis mit Dateinamen
        filename = os.path.basename(src_file_path)
        if sensitive_classification['contains_sensitive_data']:
            print(f"  ⚠️  DSGVO [{filename}]: {', '.join(sensitive_classification['data_categories'])} "
                  f"(Schutzklasse: {sensitive_classification['protection_level']})")
        else:
            print(f"  ✓ DSGVO [{filename}]: Keine besonders schutzbedürftigen Daten")

        return True

    except Exception as e:
        print(f"Fehler beim DSGVO-Update: {e}")
        import traceback
        traceback.print_exc()
        return False

def validate_json_file(json_path, src_file_path=None):
    """
    Validiert eine JSON-Ausgabedatei auf Korrektheit und sinnvollen Inhalt.

    Args:
        json_path: Pfad zur JSON-Datei
        src_file_path: Optional - Pfad zur Quelldatei für Zeitstempelprüfung

    Returns:
        True: Datei ist valide und kann übersprungen werden
        False: Datei ist fehlerhaft und muss neu erstellt werden
    """
    if not os.path.exists(json_path):
        return False

    try:
        with open(json_path, 'r', encoding='utf-8') as f:
            data = json.load(f)

        # Prüfe erforderliche Felder
        required_fields = ['path', 'ext', 'size', 'created', 'modified', 'chars', 'summary']
        for field in required_fields:
            if field not in data:
                print(f"Fehlende Struktur in {json_path}: Feld '{field}' fehlt")
                return False

        # Prüfe Content-Hash wenn Quelldatei angegeben wurde (Hash-basierte Änderungserkennung)
        if src_file_path and os.path.exists(src_file_path):
            try:
                # Wenn kein Hash in JSON vorhanden, nutze Zeitstempel (Fallback für alte JSON-Dateien)
                if 'content_hash' not in data:
                    stat = os.stat(src_file_path)
                    current_created = datetime.fromtimestamp(stat.st_ctime).isoformat()
                    current_modified = datetime.fromtimestamp(stat.st_mtime).isoformat()

                    # Vergleiche Zeitstempel
                    if data.get('created') != current_created or data.get('modified') != current_modified:
                        print(f"Zeitstempel geändert in {json_path} - Neuverarbeitung erforderlich")
                        return False
                else:
                    # Hash-basierte Prüfung (bevorzugte Methode)
                    current_hash = calculate_content_hash(src_file_path)
                    if current_hash and data.get('content_hash') != current_hash:
                        print(f"Dateiinhalt geändert in {json_path} - Neuverarbeitung erforderlich")
                        return False
            except Exception as e:
                print(f"Fehler beim Prüfen der Dateiänderungen für {src_file_path}: {e}")
                # Bei Fehler trotzdem als valide betrachten (sicherer)

        # Prüfe ob Summary sinnvoll ist (nicht leer, nicht nur Leerzeichen)
        summary = data.get('summary', '').strip()
        if not summary:
            print(f"Leere Summary in {json_path}")
            return False

        # Prüfe auf typische Fehlermarker
        error_indicators = [
            'error', 'fehler', 'failed', 'exception',
            'cannot', 'kann nicht', 'konnte nicht'
        ]
        summary_lower = summary.lower()
        if any(indicator in summary_lower for indicator in error_indicators):
            # Nur als Fehler werten, wenn die Summary sehr kurz ist (wahrscheinlich Fehlermeldung)
            if len(summary) < 50:
                print(f"Verdächtige Summary in {json_path}: {summary[:100]}")
                return False

        # Prüfe Mindestlänge der Summary
        if len(summary) < 20:
            print(f"Summary zu kurz in {json_path}: {summary}")
            return False

        # Datei scheint valide zu sein
        return True

    except json.JSONDecodeError as e:
        print(f"JSON-Fehler in {json_path}: {e}")
        return False
    except Exception as e:
        print(f"Fehler beim Validieren von {json_path}: {e}")
        return False

def check_user_input():
    """
    Prüft ob eine Taste gedrückt wurde (nicht-blockierend) - macOS Version.
    Returns True wenn eine Taste gedrückt wurde.
    """
    try:
        # Verwende select() für nicht-blockierende Eingabe auf macOS/Unix
        rlist, _, _ = select.select([sys.stdin], [], [], 0)
        if rlist:
            # Lese und verwerfe die Eingabe
            sys.stdin.readline()
            return True
        return False
    except Exception:
        # Falls select nicht funktioniert, gebe False zurück
        return False

def ask_continue():
    """
    Fragt den Benutzer höflich, ob er fortfahren möchte.
    Returns True wenn fortgesetzt werden soll, False zum Abbrechen.
    """
    print("\n" + "!" * 70)
    print("PAUSE - Eine Taste wurde gedrückt")
    print("!" * 70)
    print("\nMöchten Sie die Verarbeitung fortsetzen?")
    print("  [J] Ja, fortfahren")
    print("  [N] Nein, abbrechen und beenden")
    print("\nBitte wählen Sie (J/N): ", end="", flush=True)

    while True:
        choice = input().strip().upper()  # Konvertiere zu Großbuchstaben

        if choice == 'J':
            print("\nVerarbeitung wird fortgesetzt...\n")
            return True
        elif choice == 'N':
            print("\nVerarbeitung wird abgebrochen. Vielen Dank!\n")
            return False
        else:
            print("Ungültige Eingabe. Bitte J oder N eingeben: ", end="", flush=True)

def ask_on_lmstudio_error(error_message, file_path):
    """
    Fragt den Benutzer, wie bei LM Studio Fehlern verfahren werden soll.
    Returns: 'abort', 'skip_prompts', oder 'continue'
    """
    global ERROR_HANDLING_MODE

    # Wenn bereits eine Entscheidung getroffen wurde
    if ERROR_HANDLING_MODE == "skip":
        return "skip_prompts"
    elif ERROR_HANDLING_MODE == "ask":
        return "continue"

    # Zeige professionelle Fehlermeldung
    print("\n" + "=" * 80)
    print("FEHLER BEI DER VERARBEITUNG")
    print("=" * 80)
    print(f"\nDatei: {os.path.basename(file_path)}")
    print(f"Fehler: {error_message}")
    print("\n" + "-" * 80)
    print("\nWie möchten Sie fortfahren?")
    print("\n  [A] Abbrechen - Verarbeitung sofort beenden")
    print("  [W] Weiter ohne Fehlerabfragen - Weitere Fehler stillschweigend überspringen")
    print("  [F] Weiter mit Fehlerabfragen - Bei jedem Fehler erneut nachfragen")
    print("\n" + "-" * 80)
    print("Bitte wählen Sie (A/W/F): ", end="", flush=True)

    while True:
        choice = input().strip().upper()  # Konvertiere zu Großbuchstaben

        if choice == 'A':
            print("\n→ Verarbeitung wird abgebrochen.\n")
            print("=" * 80 + "\n")
            return "abort"
        elif choice == 'W':
            print("\n→ Verarbeitung wird fortgesetzt. Weitere Fehler werden nicht mehr angezeigt.\n")
            print("=" * 80 + "\n")
            ERROR_HANDLING_MODE = "skip"
            return "skip_prompts"
        elif choice == 'F':
            print("\n→ Verarbeitung wird fortgesetzt. Bei Fehlern erfolgt eine erneute Abfrage.\n")
            print("=" * 80 + "\n")
            ERROR_HANDLING_MODE = "ask"
            return "continue"
        else:
            print("Ungültige Eingabe. Bitte A, W oder F eingeben: ", end="", flush=True)

def check_lmstudio_connection():
    """
    Prüft ob LM Studio läuft und erreichbar ist.
    Returns True wenn verbunden, False sonst.
    """
    try:
        # Versuche eine einfache Anfrage an den Health-Endpoint
        health_url = LMSTUDIO_API_URL.replace('/v1/chat/completions', '/v1/models')
        response = requests.get(health_url, timeout=5)
        return response.status_code == 200
    except requests.exceptions.RequestException:
        return False

def check_ocr_functionality():
    """
    Prüft ob OCR (Tesseract) korrekt installiert und funktionsfähig ist.
    Returns: (is_available, error_message)
    """
    if not OCR_AVAILABLE:
        return False, "pytesseract oder Pillow nicht installiert"

    try:
        # Versuche tesseract Version zu prüfen
        version = pytesseract.get_tesseract_version()

        # Prüfe ob deutsche Sprache verfügbar ist
        try:
            langs = pytesseract.get_languages()
            if 'deu' not in langs:
                return True, f"Tesseract {version} verfügbar, aber deutsche Sprache 'deu' fehlt"
        except:
            # Wenn get_languages fehlschlägt, gehen wir davon aus dass es funktioniert
            pass

        return True, f"Tesseract {version} mit deutscher Sprache verfügbar"

    except pytesseract.TesseractNotFoundError:
        return False, "Tesseract Binary nicht gefunden (nicht installiert oder nicht im PATH)"
    except Exception as e:
        return False, f"OCR-Fehler: {str(e)}"

def format_time(seconds):
    """Formatiert Sekunden in h:mm:ss Format."""
    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    secs = int(seconds % 60)

    if hours > 0:
        return f"{hours}h {minutes:02d}m {secs:02d}s"
    elif minutes > 0:
        return f"{minutes}m {secs:02d}s"
    else:
        return f"{secs}s"

def walk_and_process():
    global ERROR_HANDLING_MODE

    # Setze Fehlerbehandlungsmodus zurück für neuen Durchlauf
    ERROR_HANDLING_MODE = None

    # Zeige Version und Startinformationen
    print("=" * 70)
    print(f"{SCRIPT_NAME}")
    print(f"Version {VERSION} vom {VERSION_DATE}")
    print("=" * 70)
    print(f"Quellverzeichnis: {SRC_ROOT}")
    print(f"Zielverzeichnis:  {DST_ROOT}")
    print(f"Dateitypen:       {', '.join(EXTENSIONS)}")
    print("=" * 70)

    # Prüfe LM Studio Verbindung
    print("\nPrüfe LM Studio Verbindung...")
    if not check_lmstudio_connection():
        print("\n" + "!" * 70)
        print("FEHLER: LM Studio ist nicht erreichbar!")
        print("!" * 70)
        print(f"\nBitte stellen Sie sicher, dass:")
        print(f"  1. LM Studio gestartet ist")
        print(f"  2. Ein Modell geladen ist")
        print(f"  3. Der Local Server läuft")
        print(f"  4. Die URL korrekt ist: {LMSTUDIO_API_URL}")
        print("\nProgramm wird beendet.")
        print("=" * 70)
        return
    print("✓ LM Studio verbunden")

    # Prüfe OCR-Funktionalität
    print("\nPrüfe OCR-Funktionalität (Tesseract)...")
    ocr_ok, ocr_message = check_ocr_functionality()
    if ocr_ok:
        print(f"✓ {ocr_message}")
    else:
        print(f"⚠ OCR nicht verfügbar: {ocr_message}")
        print("  Hinweis: Gescannte PDFs werden übersprungen.")
        print("  Installation:")
        print("    macOS:  brew install tesseract tesseract-lang")
        print("    Linux:  sudo apt-get install tesseract-ocr tesseract-ocr-deu")
        print("    Python: pip install pytesseract pillow")

    # Zähle zunächst alle zu verarbeitenden Dateien mit Fortschrittsanzeige
    print("\nScanne Verzeichnis...")
    all_files = []
    dir_count = 0
    file_count = 0
    last_line_length = 0  # Tracke die Länge der letzten Zeile
    file_stats = {}  # Statistik: {extension: {'count': n, 'size': bytes, 'files': [paths]}}

    # Nutze os.walk und zeige Fortschritt mit \r (carriage return)
    for root, dirs, files in os.walk(SRC_ROOT):
        # Sortiere Verzeichnisse und Dateien alphabetisch
        dirs.sort()
        files.sort()

        dir_count += 1

        # Zähle relevante Dateien in diesem Verzeichnis
        relevant_files = 0
        for name in files:
            file_count += 1
            ext = os.path.splitext(name)[1].lower()
            full_path = os.path.join(root, name)

            # Sammle Statistiken für alle Dateien
            try:
                file_size = os.path.getsize(full_path)
                if ext not in file_stats:
                    file_stats[ext] = {'count': 0, 'size': 0, 'files': []}
                file_stats[ext]['count'] += 1
                file_stats[ext]['size'] += file_size
                file_stats[ext]['files'].append(full_path)
            except (OSError, PermissionError):
                # Überspringe Dateien, auf die nicht zugegriffen werden kann
                pass

            if ext in EXTENSIONS:
                all_files.append(full_path)
                relevant_files += 1

        # Zeige Fortschritt (überschreibe vorherige Zeile mit \r)
        if dir_count % 10 == 0 or relevant_files > 0:
            rel_path = os.path.relpath(root, SRC_ROOT)
            if len(rel_path) > 50:
                rel_path = "..." + rel_path[-47:]

            # Erstelle die Ausgabezeile
            line = f"  Verzeichnisse: {dir_count:,} | Dateien gescannt: {file_count:,} | Relevante Dateien: {len(all_files):,} | {rel_path}"

            # Füge Leerzeichen hinzu, um alte längere Zeilen vollständig zu überschreiben
            if len(line) < last_line_length:
                line = line + " " * (last_line_length - len(line))

            # Merke die aktuelle Zeilenlänge für das nächste Update
            last_line_length = len(line)

            # Schreibe mit \r am Anfang (Cursor an Zeilenanfang) ohne Zeilenumbruch
            sys.stdout.write(f"\r{line}")
            sys.stdout.flush()

    # Finaler Status mit Zeilenumbruch
    final_line = f"  Verzeichnisse: {dir_count:,} | Dateien gescannt: {file_count:,} | Relevante Dateien: {len(all_files):,}"
    if len(final_line) < last_line_length:
        final_line = final_line + " " * (last_line_length - len(final_line))
    sys.stdout.write(f"\r{final_line}\n")
    sys.stdout.flush()

    # Zeige Statistik nach Dateiendungen
    print("\n" + "=" * 80)
    print("STATISTIK DER DATEIENDUNGEN")
    print("=" * 80)
    print(f"{'Endung':<18} {'Anzahl':>8} {'Größe (MB)':>12} {'Ø Größe (KB)':>14}  {'Status':<15}")
    print("-" * 80)

    # Sortiere nach Anzahl der Dateien (absteigend)
    sorted_stats = sorted(file_stats.items(), key=lambda x: x[1]['count'], reverse=True)

    for ext, stats in sorted_stats:
        count = stats['count']
        total_size_mb = stats['size'] / (1024 * 1024)
        avg_size_kb = (stats['size'] / count) / 1024 if count > 0 else 0

        # Markiere ob dieser Typ analysiert wird
        ext_display = ext if ext else "(keine)"
        will_analyze = "→ WIRD ANALYSIERT" if ext in EXTENSIONS else ""

        print(f"{ext_display:<18} {count:>8,} {total_size_mb:>12.2f} {avg_size_kb:>14.2f}  {will_analyze:<15}")

    print("=" * 80)

    total_files = len(all_files)
    print(f"\nGefunden: {total_files} Dateien zum Verarbeiten")
    print("\nHinweis: Drücken Sie Enter während der Verarbeitung,")
    print("         um anzuhalten und zu wählen, ob Sie fortfahren möchten.")
    print("=" * 70)

    if total_files == 0:
        print("Keine Dateien gefunden.")
        return

    # Verarbeite Dateien mit Fortschrittsanzeige
    processed = 0
    skipped = 0
    errors = 0
    recreated = 0
    ocr_count = 0  # Zähler für OCR-verarbeitete Dokumente
    excluded = 0   # Zähler für ausgeschlossene Verzeichnisse
    duplicates = 0  # Zähler für Duplikate
    start_time = time.time()

    for idx, full_path in enumerate(all_files, 1):
        # Prüfe auf Tasteneingabe
        if check_user_input():
            if not ask_continue():
                print("\n" + "=" * 70)
                print("VERARBEITUNG VOM BENUTZER ABGEBROCHEN")
                print("=" * 70)
                print(f"Verarbeitet bis Datei {idx}/{total_files}")
                print(f"Neu verarbeitet: {processed}")
                print(f"Neu erstellt (vorher fehlerhaft): {recreated}")
                print(f"Übersprungen (valide): {skipped}")
                print(f"Fehler: {errors}")
                print(f"Mit OCR verarbeitet: {ocr_count}")
                elapsed = time.time() - start_time
                print(f"Laufzeit bis Abbruch: {format_time(elapsed)}")
                print("=" * 70)
                return

        try:
            # Schritt 1: Prüfe ob Pfad ausgeschlossen werden soll
            if should_exclude_path(full_path):
                excluded += 1
                if excluded <= 10:  # Zeige nur erste 10
                    print(f"Ausgeschlossen (Pattern-Match): {os.path.relpath(full_path, SRC_ROOT)}")
                continue

            # Schritt 2: Prüfe auf Duplikate (basierend auf Content-Hash)
            try:
                file_size = os.path.getsize(full_path)
                is_dup, original_path = is_duplicate_file(full_path, file_size)
                if is_dup:
                    duplicates += 1
                    if duplicates <= 10:  # Zeige nur erste 10
                        print(f"Duplikat übersprungen: {os.path.relpath(full_path, SRC_ROOT)}")
                        print(f"  → Original: {os.path.relpath(original_path, SRC_ROOT)}")
                    continue
            except OSError:
                pass  # Bei Fehler: Fahre normal fort

            # Schritt 3: Prüfe ob bereits existiert und valide ist
            rel_path = os.path.relpath(full_path, SRC_ROOT)
            dst_dir = os.path.join(DST_ROOT, os.path.dirname(rel_path))
            dst_file = os.path.join(dst_dir, os.path.basename(full_path) + ".json")

            ocr_info = None
            if os.path.exists(dst_file):
                if validate_json_file(dst_file, full_path):
                    # JSON ist valide - prüfe ob Kontaktinformationen nachgetragen werden müssen
                    updated = update_json_with_contact_info(dst_file, full_path)
                    if not updated:
                        skipped += 1
                    else:
                        skipped += 1  # Zählt trotzdem als übersprungen (nur Mini-Update)

                    # Lese OCR-Info aus existierender JSON-Datei für Statistik
                    try:
                        with open(dst_file, 'r', encoding='utf-8') as f:
                            existing_data = json.load(f)
                            ocr_info = existing_data.get('ocr_info', None)
                    except:
                        pass
                else:
                    # Fehlerhafte oder veraltete Datei wird in process_file gelöscht und neu erstellt
                    ocr_info = process_file(full_path)
                    recreated += 1
            else:
                ocr_info = process_file(full_path)
                processed += 1

            # Zähle OCR-verarbeitete Dokumente
            if ocr_info and ocr_info.get('used_ocr'):
                ocr_count += 1

            # Berechne Zeitschätzung
            elapsed = time.time() - start_time
            if idx > 0:
                # Berechne Durchschnitt nur für tatsächlich verarbeitete Dateien
                actually_processed = processed + recreated
                # Zeige aktuellen Dateipfad (relativ für bessere Lesbarkeit)
                current_rel_path = os.path.relpath(full_path, SRC_ROOT)
                print(f"\n📄 Datei: {current_rel_path}")

                if actually_processed > 0:
                    avg_time_per_file = elapsed / actually_processed
                    remaining_files = total_files - idx
                    estimated_remaining = avg_time_per_file * remaining_files

                    print(f"[{idx}/{total_files}] Fortschritt: {(idx/total_files)*100:.1f}%")
                    print(f"Neu: {processed} | Neu erstellt: {recreated} | Übersprungen: {skipped} | Fehler: {errors}")
                    print(f"Duplikate: {duplicates} | Ausgeschlossen: {excluded} | OCR: {ocr_count}")
                    print(f"Verstrichene Zeit: {format_time(elapsed)}")
                    print(f"Geschätzte Restzeit: {format_time(estimated_remaining)}")
                    print(f"Geschätzte Gesamtzeit: {format_time(elapsed + estimated_remaining)}")
                    print(f"Durchschnitt: {avg_time_per_file:.2f}s pro Datei")
                    print("=" * 70)
                else:
                    print(f"[{idx}/{total_files}] Fortschritt: {(idx/total_files)*100:.1f}%")
                    print(f"Neu: {processed} | Neu erstellt: {recreated} | Übersprungen: {skipped} | Fehler: {errors}")
                    print(f"Duplikate: {duplicates} | Ausgeschlossen: {excluded} | OCR: {ocr_count}")
                    print("=" * 70)

        except Exception as e:
            errors += 1
            print("Fehler bei", full_path, "->", e)

    # Abschlussbericht
    total_time = time.time() - start_time
    print("\n" + "=" * 70)
    print("VERARBEITUNG ABGESCHLOSSEN")
    print("=" * 70)
    print(f"Gesamt gescannt: {total_files} Dateien")
    print(f"Neu verarbeitet: {processed}")
    print(f"Neu erstellt (vorher fehlerhaft): {recreated}")
    print(f"Übersprungen (valide): {skipped}")
    print(f"Duplikate übersprungen: {duplicates}")
    print(f"Ausgeschlossen (Pattern): {excluded}")
    print(f"Fehler: {errors}")
    print(f"Mit OCR verarbeitet: {ocr_count}")
    print(f"Gesamtzeit: {format_time(total_time)}")
    # Berechne Durchschnitt nur für tatsächlich verarbeitete Dateien (nicht übersprungene)
    actually_processed = processed + recreated
    if actually_processed > 0:
        print(f"Durchschnitt: {total_time/actually_processed:.2f}s pro Datei (nur verarbeitete)")
    if duplicates > 0:
        print(f"\nℹ Hinweis: {duplicates} Duplikate wurden automatisch erkannt und übersprungen")
    if excluded > 0:
        print(f"ℹ Hinweis: {excluded} Dateien in ausgeschlossenen Verzeichnissen übersprungen")
    print("=" * 70)

def cleanup_invalid_phone_numbers():
    """
    Bereinigt alle vorhandenen JSON-Dateien und entfernt ungültige Kontaktinformationen.
    Re-extrahiert URLs, E-Mails und Telefonnummern aus Quelldateien wenn nötig.
    Entfernt: URLs mit Satzzeichen, E-Mails mit URL-Präfix, ungültige Telefonnummern.
    """
    print("\n" + "=" * 80)
    print("BEREINIGUNG: UNGÜLTIGE KONTAKTINFORMATIONEN ENTFERNEN")
    print("=" * 80)
    print(f"Durchsuche: {DST_ROOT}")
    print("=" * 80 + "\n")

    # Sammle alle JSON-Dateien
    all_json_files = []
    for root, dirs, files in os.walk(DST_ROOT):
        # Sortiere für konsistente Reihenfolge
        dirs.sort()
        files.sort()

        for name in files:
            if name.endswith('.json'):
                full_path = os.path.join(root, name)
                all_json_files.append(full_path)

    total_files = len(all_json_files)
    print(f"Gefunden: {total_files:,} JSON-Dateien\n")

    if total_files == 0:
        print("Keine JSON-Dateien gefunden.")
        return

    # Statistiken
    files_cleaned = 0
    total_invalid_removed = 0
    files_with_invalid = 0
    start_time = time.time()

    for idx, json_file in enumerate(all_json_files, 1):
        try:
            # Bestimme Quelldatei
            # JSON-Dateien enden mit ".original_extension.json"
            rel_path = os.path.relpath(json_file, DST_ROOT)
            src_rel_path = rel_path.replace('.json', '')  # Entferne .json
            src_file = os.path.join(SRC_ROOT, src_rel_path)

            if not os.path.exists(src_file):
                continue

            # Rufe update_json_with_contact_info auf (prüft und bereinigt automatisch)
            was_updated = update_json_with_contact_info(json_file, src_file)

            if was_updated:
                files_cleaned += 1

            # Fortschritt anzeigen
            if idx % 100 == 0 or idx == total_files:
                elapsed = time.time() - start_time
                progress = (idx / total_files) * 100
                print(f"\rFortschritt: {idx:,}/{total_files:,} ({progress:.1f}%) - "
                      f"Bereinigt: {files_cleaned:,} - Zeit: {elapsed:.1f}s", end="", flush=True)

        except Exception as e:
            print(f"\nFehler bei {json_file}: {e}")

    # Abschlussbericht
    total_time = time.time() - start_time
    print(f"\n\n" + "=" * 80)
    print("BEREINIGUNG ABGESCHLOSSEN")
    print("=" * 80)
    print(f"Gescannte Dateien: {total_files:,}")
    print(f"Bereinigte Dateien: {files_cleaned:,}")
    print(f"Gesamtzeit: {format_time(total_time)}")
    print("=" * 80)

def update_all_jsons_with_dsgvo():
    """
    Aktualisiert alle vorhandenen JSON-Dateien mit DSGVO-Klassifizierung.
    Analysiert Dokumente auf besonders schutzbedürftige personenbezogene Daten
    gemäß Art. 9 DSGVO und § 26 BDSG.
    Sehr schnell - nur Regex, kein LLM.
    """
    print("\n" + "=" * 80)
    print("DSGVO-UPDATE: KLASSIFIZIERUNG BESONDERS SCHUTZBEDÜRFTIGER DATEN")
    print("=" * 80)
    print(f"Analysiere Dokumente gemäß Art. 9 DSGVO und § 26 BDSG")
    print(f"Durchsuche: {DST_ROOT}")
    print("=" * 80 + "\n")

    # Sammle alle JSON-Dateien
    all_json_files = []
    for root, dirs, files in os.walk(DST_ROOT):
        # Sortiere für konsistente Reihenfolge
        dirs.sort()
        files.sort()

        for name in files:
            if name.endswith('.json'):
                full_path = os.path.join(root, name)
                all_json_files.append(full_path)

    total_files = len(all_json_files)
    print(f"Gefunden: {total_files:,} JSON-Dateien\n")

    if total_files == 0:
        print("Keine JSON-Dateien gefunden.")
        return

    # Statistiken
    files_updated = 0
    files_with_sensitive_data = 0
    sensitive_categories = {}
    start_time = time.time()

    for idx, json_file in enumerate(all_json_files, 1):
        try:
            # Bestimme Quelldatei
            # JSON-Dateien enden mit ".original_extension.json"
            rel_path = os.path.relpath(json_file, DST_ROOT)
            src_rel_path = rel_path.replace('.json', '')  # Entferne .json
            src_file = os.path.join(SRC_ROOT, src_rel_path)

            if not os.path.exists(src_file):
                continue

            # Rufe DSGVO-Update auf
            was_updated = update_json_with_dsgvo_classification(json_file, src_file)

            if was_updated:
                files_updated += 1

                # Lese JSON um zu prüfen ob sensible Daten gefunden wurden
                try:
                    with open(json_file, 'r', encoding='utf-8') as f:
                        data = json.load(f)
                        dsgvo = data.get('dsgvo_classification', {})
                        if dsgvo.get('contains_sensitive_data'):
                            files_with_sensitive_data += 1
                            # Sammle Kategorien für Statistik
                            for category in dsgvo.get('data_categories', []):
                                sensitive_categories[category] = sensitive_categories.get(category, 0) + 1
                except:
                    pass

            # Fortschritt anzeigen
            if idx % 50 == 0 or idx == total_files:
                elapsed = time.time() - start_time
                progress = (idx / total_files) * 100
                print(f"\rFortschritt: {idx:,}/{total_files:,} ({progress:.1f}%) - "
                      f"Aktualisiert: {files_updated:,} - Sensible Daten: {files_with_sensitive_data:,} - "
                      f"Zeit: {elapsed:.1f}s", end="", flush=True)

        except Exception as e:
            print(f"\nFehler bei {json_file}: {e}")

    # Abschlussbericht
    total_time = time.time() - start_time
    print(f"\n\n" + "=" * 80)
    print("DSGVO-UPDATE ABGESCHLOSSEN")
    print("=" * 80)
    print(f"Gescannte Dateien: {total_files:,}")
    print(f"Aktualisierte Dateien: {files_updated:,}")
    print(f"Dateien mit sensiblen Daten: {files_with_sensitive_data:,}")

    if sensitive_categories:
        print(f"\nGefundene Kategorien besonders schutzbedürftiger Daten:")
        for category, count in sorted(sensitive_categories.items(), key=lambda x: x[1], reverse=True):
            print(f"  • {category}: {count:,} Dokumente")

    print(f"\nGesamtzeit: {format_time(total_time)}")
    print("=" * 80)

def create_combined_database(max_size_mb=30, output_dir=None):
    """
    Erstellt kombinierte JSON-Datenbank-Dateien aus allen einzelnen JSON-Dateien.
    Teilt die Datenbank in mehrere Dateien auf, wenn die Größe max_size_mb überschreitet.

    Args:
        max_size_mb: Maximale Größe pro Datenbankdatei in MB
        output_dir: Ausgabeverzeichnis für Datenbankdateien (Standard: DST_ROOT/database)
    """
    if output_dir is None:
        output_dir = os.path.join(DST_ROOT, "database")

    # Erstelle Ausgabeverzeichnis
    os.makedirs(output_dir, exist_ok=True)

    print("\n" + "=" * 80)
    print("ERSTELLE KOMBINIERTE JSON-DATENBANK")
    print("=" * 80)
    print(f"Quellverzeichnis: {DST_ROOT}")
    print(f"Ausgabeverzeichnis: {output_dir}")
    print(f"Maximale Größe pro Datei: {max_size_mb} MB")
    print("=" * 80)

    # Sammle alle JSON-Dateien
    print("\nSammle JSON-Dateien...")
    all_json_files = []
    for root, dirs, files in os.walk(DST_ROOT):
        # Überspringe das database-Verzeichnis selbst
        if root.startswith(output_dir):
            continue

        # Sortiere für konsistente Reihenfolge
        dirs.sort()
        files.sort()

        for name in files:
            if name.endswith('.json'):
                full_path = os.path.join(root, name)
                all_json_files.append(full_path)

    total_files = len(all_json_files)
    print(f"Gefunden: {total_files:,} JSON-Dateien")

    if total_files == 0:
        print("Keine JSON-Dateien gefunden. Bitte führen Sie zuerst die normale Verarbeitung durch.")
        return

    # Lade und kombiniere JSON-Dateien
    print("\nLade und kombiniere Dateien...")
    max_size_bytes = max_size_mb * 1024 * 1024

    current_batch = []
    current_size = 0
    batch_number = 1
    total_size = 0
    failed_files = 0

    # Metadaten für die Datenbank
    database_metadata = {
        "created": datetime.now().isoformat(),
        "source_directory": SRC_ROOT,
        "json_directory": DST_ROOT,
        "total_documents": 0,
        "script_version": VERSION,
        "script_date": VERSION_DATE
    }

    start_time = time.time()

    for idx, json_file in enumerate(all_json_files, 1):
        try:
            with open(json_file, 'r', encoding='utf-8') as f:
                data = json.load(f)

            # Schätze die Größe dieses Eintrags
            entry_json = json.dumps(data, ensure_ascii=False)
            entry_size = len(entry_json.encode('utf-8'))

            # Prüfe ob wir eine neue Datei starten müssen
            # Reserviere 2000 Bytes für Metadaten und JSON-Struktur
            if current_size + entry_size + 2000 > max_size_bytes and current_batch:
                # Schreibe aktuelle Batch
                write_database_file(output_dir, batch_number, current_batch, database_metadata, max_size_mb)
                total_size += current_size
                batch_number += 1
                current_batch = []
                current_size = 0

            # Füge zu aktueller Batch hinzu
            current_batch.append(data)
            current_size += entry_size

            # Fortschrittsanzeige
            if idx % 100 == 0 or idx == total_files:
                progress = (idx / total_files) * 100
                print(f"\rFortschritt: {idx:,}/{total_files:,} ({progress:.1f}%) - "
                      f"Batch {batch_number}: {len(current_batch):,} Dateien, "
                      f"{current_size / (1024*1024):.2f} MB", end="", flush=True)

        except json.JSONDecodeError as e:
            failed_files += 1
            print(f"\nWarnung: Fehlerhafte JSON-Datei übersprungen: {json_file}")
            print(f"  Fehler: {e}")
        except Exception as e:
            failed_files += 1
            print(f"\nWarnung: Fehler beim Lesen von {json_file}: {e}")

    print()  # Neue Zeile nach Fortschrittsanzeige

    # Schreibe letzte Batch
    if current_batch:
        write_database_file(output_dir, batch_number, current_batch, database_metadata, max_size_mb)
        total_size += current_size

    # Abschlussbericht
    elapsed = time.time() - start_time
    total_documents = sum(len(current_batch) if i == batch_number else 0 for i in range(1, batch_number + 1))

    # Berechne korrekte Gesamtanzahl Dokumente
    total_documents = total_files - failed_files

    print("\n" + "=" * 80)
    print("DATENBANK-ERSTELLUNG ABGESCHLOSSEN")
    print("=" * 80)
    print(f"Verarbeitete Dokumente: {total_documents:,}")
    print(f"Fehlerhafte Dateien: {failed_files:,}")
    print(f"Anzahl Datenbank-Dateien: {batch_number}")
    print(f"Gesamtgröße: {total_size / (1024*1024):.2f} MB")
    print(f"Durchschnittliche Größe pro Datei: {(total_size / batch_number) / (1024*1024):.2f} MB")
    print(f"Ausgabeverzeichnis: {output_dir}")
    print(f"Laufzeit: {format_time(elapsed)}")
    print("=" * 80)

    # Liste der erstellten Dateien
    print("\nErstellte Datenbank-Dateien:")
    for i in range(1, batch_number + 1):
        filename = f"file_database_{i:03d}.json"
        filepath = os.path.join(output_dir, filename)
        if os.path.exists(filepath):
            size_mb = os.path.getsize(filepath) / (1024 * 1024)
            print(f"  {filename}: {size_mb:.2f} MB")
    print("=" * 80)

def write_database_file(output_dir, batch_number, documents, metadata, max_size_mb):
    """
    Schreibt eine Datenbank-Datei mit Metadaten und Dokumenten.

    Args:
        output_dir: Ausgabeverzeichnis
        batch_number: Nummer der Batch (für Dateinamen)
        documents: Liste der Dokumente
        metadata: Metadaten für die Datenbank
        max_size_mb: Maximale Größe (für Metadaten)
    """
    filename = f"file_database_{batch_number:03d}.json"
    filepath = os.path.join(output_dir, filename)

    # Erweitere Metadaten
    batch_metadata = metadata.copy()
    batch_metadata["batch_number"] = batch_number
    batch_metadata["documents_in_batch"] = len(documents)
    batch_metadata["max_size_mb"] = max_size_mb

    # Erstelle Datenbank-Struktur
    database = {
        "metadata": batch_metadata,
        "documents": documents
    }

    # Schreibe Datei
    with open(filepath, 'w', encoding='utf-8') as f:
        json.dump(database, f, ensure_ascii=False, indent=2)

    size_mb = os.path.getsize(filepath) / (1024 * 1024)
    print(f"\n✓ Erstellt: {filename} ({size_mb:.2f} MB, {len(documents):,} Dokumente)")

def parse_arguments():
    """Parse und validiere Kommandozeilenargumente."""
    parser = argparse.ArgumentParser(
        description=f'{SCRIPT_NAME} - Version {VERSION}',
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=f"""
Beispiele:
  {sys.argv[0]}
    Verwendet Standard-Verzeichnisse und -Einstellungen

  {sys.argv[0]} --src ~/Documents --dst ~/Summaries
    Verwendet benutzerdefinierte Verzeichnisse

  {sys.argv[0]} --max-tokens 8192
    Verwendet kleineres Modell mit 8k Tokens Context

  {sys.argv[0]} --summary-max-chars 2000
    Erstellt längere Zusammenfassungen (max 2000 Zeichen)

  {sys.argv[0]} --src ~/Docs --dst ~/Summaries --max-tokens 32768 --summary-max-chars 2000
    Vollständig benutzerdefinierte Konfiguration

  {sys.argv[0]} --create-database
    Erstellt kombinierte Datenbank aus allen JSON-Dateien (Standard: max 30 MB pro Datei)

  {sys.argv[0]} --create-database --max-database-size 50
    Erstellt Datenbank mit max 50 MB pro Datei

  {sys.argv[0]} --create-database --database-output ~/MyDatabase
    Erstellt Datenbank in benutzerdefiniertem Verzeichnis

  {sys.argv[0]} --cleanup-phones
    Bereinigt alle JSON-Dateien: Entfernt ungültige Telefonnummern (z.B. Projektnummern)
    und extrahiert korrekte Telefonnummern neu aus den Quelldateien

  {sys.argv[0]} --update-dsgvo
    Aktualisiert alle bestehenden JSON-Dateien mit DSGVO-Klassifizierung.
    Analysiert Dokumente auf besonders schutzbedürftige personenbezogene Daten
    gemäß Art. 9 DSGVO und § 26 BDSG (sehr schnell, kein LLM)

  {sys.argv[0]} --version
    Zeigt Versionsinformation an

Konfiguration:
  Die Standardwerte können in der Datei direkt angepasst werden:
    SRC_ROOT = "~/OneDrive - CompanyName"
    DST_ROOT = "~/LLM"
    MAX_CONTEXT_TOKENS = 262144

Empfohlene MAX_CONTEXT_TOKENS Werte:
  - Kleinere Modelle (z.B. Llama 3 8B): 8192
  - Größere Modelle (z.B. Qwen 2.5 14B): 32768
  - Reasoning-Modelle (z.B. ministral-3-14b-reasoning): 262144

Weitere Informationen:
  Siehe README.md für detaillierte Dokumentation
        """
    )

    parser.add_argument(
        '--src',
        type=str,
        help=f'Quellverzeichnis (Standard: {SRC_ROOT})'
    )

    parser.add_argument(
        '--dst',
        type=str,
        help=f'Zielverzeichnis für JSON-Dateien (Standard: {DST_ROOT})'
    )

    parser.add_argument(
        '--max-tokens',
        type=int,
        metavar='TOKENS',
        help=f'Maximale Context-Länge des Modells in Tokens (Standard: {MAX_CONTEXT_TOKENS})'
    )

    parser.add_argument(
        '--summary-max-chars',
        type=int,
        metavar='CHARS',
        help=f'Maximale Länge der Zusammenfassung in Zeichen. Text kürzer als dieser Wert wird direkt kopiert. (Standard: {SUMMARY_MAX_CHARS})'
    )

    parser.add_argument(
        '--version',
        action='version',
        version=f'{SCRIPT_NAME}\nVersion: {VERSION}\nDatum: {VERSION_DATE}'
    )

    parser.add_argument(
        '--create-database',
        action='store_true',
        help='Erstellt eine kombinierte JSON-Datenbank aus allen einzelnen JSON-Dateien (max. 30 MB pro Datei)'
    )

    parser.add_argument(
        '--database-output',
        type=str,
        metavar='DIR',
        help='Ausgabeverzeichnis für die Datenbank-Dateien (Standard: DST_ROOT/database)'
    )

    parser.add_argument(
        '--cleanup-phones',
        action='store_true',
        help='Bereinigt alle JSON-Dateien: Entfernt ungültige Telefonnummern und re-extrahiert korrekte aus Quelldateien'
    )

    parser.add_argument(
        '--update-dsgvo',
        action='store_true',
        help='Aktualisiert alle JSON-Dateien mit DSGVO-Klassifizierung (Art. 9 DSGVO, § 26 BDSG)'
    )

    parser.add_argument(
        '--max-database-size',
        type=int,
        metavar='MB',
        default=30,
        help='Maximale Größe pro Datenbank-Datei in MB (Standard: 30)'
    )

    return parser.parse_args()

if __name__ == "__main__":
    # Parse Kommandozeilenargumente
    args = parse_arguments()

    # Überschreibe globale Variablen falls Parameter angegeben wurden
    if args.src:
        SRC_ROOT = os.path.expanduser(args.src)
        # Aktualisiere die globale Variable
        globals()['SRC_ROOT'] = SRC_ROOT
    if args.dst:
        DST_ROOT = os.path.expanduser(args.dst)
        # Aktualisiere die globale Variable
        globals()['DST_ROOT'] = DST_ROOT
    if args.max_tokens:
        MAX_CONTEXT_TOKENS = args.max_tokens
        # Aktualisiere die globale Variable
        globals()['MAX_CONTEXT_TOKENS'] = MAX_CONTEXT_TOKENS
    if args.summary_max_chars:
        SUMMARY_MAX_CHARS = args.summary_max_chars
        # Aktualisiere die globale Variable
        globals()['SUMMARY_MAX_CHARS'] = SUMMARY_MAX_CHARS

    # Prüfe ob Telefonnummern-Bereinigung gewünscht ist
    if args.cleanup_phones:
        cleanup_invalid_phone_numbers()
        sys.exit(0)

    # Prüfe ob DSGVO-Update gewünscht ist
    if args.update_dsgvo:
        update_all_jsons_with_dsgvo()
        sys.exit(0)

    # Prüfe ob Datenbank-Erstellung gewünscht ist
    if args.create_database:
        # Erstelle kombinierte Datenbank
        output_dir = args.database_output if args.database_output else None
        if output_dir:
            output_dir = os.path.expanduser(output_dir)
        create_combined_database(
            max_size_mb=args.max_database_size,
            output_dir=output_dir
        )
    else:
        # Normale Verarbeitung
        walk_and_process()
